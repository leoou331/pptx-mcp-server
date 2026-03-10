"""
PPTX MCP Server v3.0
基于 python-pptx 的 PowerPoint 操作服务

功能：
- 创建/打开/保存 PPTX 文件
- 添加幻灯片、文本、图片、表格
- 安全验证（ZIP炸弹、宏检测、路径遍历）
- 会话管理和超时清理

启动方式：
  python server.py --port 8010 --token your-token
"""
import asyncio
import json
import logging
import os
import signal
import sys
import argparse
from http.server import HTTPServer, BaseHTTPRequestHandler
from socketserver import ThreadingMixIn
from typing import Optional, Dict, Any

# ===== 安全配置：必须在导入 pptx 之前 =====
from lxml import etree
import lxml.etree as lxml_etree

# 保存原始构造函数
_original_XMLParser = etree.XMLParser

def _secure_XMLParser(*args, **kwargs):
    """强制安全配置的 XML 解析器"""
    kwargs.setdefault('resolve_entities', False)
    kwargs.setdefault('no_network', True)
    kwargs.setdefault('huge_tree', False)
    kwargs.setdefault('load_dtd', False)
    return _original_XMLParser(*args, **kwargs)

# Monkey Patch
etree.XMLParser = _secure_XMLParser
lxml_etree.XMLParser = _secure_XMLParser

# ===== 导入 pptx =====
from pptx import Presentation

# ===== 导入安全模块 =====
from security import (
    validate_pptx,
    has_macro,
    safe_path,
    limits,
    SessionManager,
    temp_manager
)
from tools import PptxTools

# ===== 配置 =====
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(name)s: %(message)s"
)
log = logging.getLogger("pptx-server")

# ===== MCP 工具定义 =====
TOOLS = [
    {
        "name": "pptx_create",
        "description": "创建空白 PowerPoint 演示文稿",
        "inputSchema": {
            "type": "object",
            "properties": {
                "name": {"type": "string", "description": "演示文稿名称"}
            },
            "required": ["name"]
        }
    },
    {
        "name": "pptx_open",
        "description": "打开现有的 PPTX 文件（带安全验证）",
        "inputSchema": {
            "type": "object",
            "properties": {
                "file_path": {"type": "string", "description": "文件路径"}
            },
            "required": ["file_path"]
        }
    },
    {
        "name": "pptx_save",
        "description": "保存演示文稿到文件",
        "inputSchema": {
            "type": "object",
            "properties": {
                "session_id": {"type": "string", "description": "会话 ID"},
                "output_path": {"type": "string", "description": "输出路径（可选）"}
            },
            "required": ["session_id"]
        }
    },
    {
        "name": "pptx_close",
        "description": "关闭会话释放资源",
        "inputSchema": {
            "type": "object",
            "properties": {
                "session_id": {"type": "string", "description": "会话 ID"}
            },
            "required": ["session_id"]
        }
    },
    {
        "name": "pptx_info",
        "description": "获取演示文稿详细信息",
        "inputSchema": {
            "type": "object",
            "properties": {
                "session_id": {"type": "string", "description": "会话 ID"}
            },
            "required": ["session_id"]
        }
    },
    {
        "name": "pptx_add_slide",
        "description": "添加新幻灯片",
        "inputSchema": {
            "type": "object",
            "properties": {
                "session_id": {"type": "string"},
                "layout_index": {"type": "integer", "default": 0, "description": "布局索引"}
            },
            "required": ["session_id"]
        }
    },
    {
        "name": "pptx_add_text",
        "description": "添加文本框到幻灯片",
        "inputSchema": {
            "type": "object",
            "properties": {
                "session_id": {"type": "string"},
                "slide_index": {"type": "integer", "description": "幻灯片索引"},
                "text": {"type": "string", "description": "文本内容"},
                "position": {"type": "string", "enum": ["title", "body", "custom"], "default": "body"},
                "left": {"type": "number", "default": 1},
                "top": {"type": "number", "default": 1},
                "width": {"type": "number", "default": 8},
                "height": {"type": "number", "default": 1},
                "font_size": {"type": "integer", "default": 18}
            },
            "required": ["session_id", "slide_index", "text"]
        }
    },
    {
        "name": "pptx_add_image",
        "description": "添加图片到幻灯片",
        "inputSchema": {
            "type": "object",
            "properties": {
                "session_id": {"type": "string"},
                "slide_index": {"type": "integer"},
                "image_path": {"type": "string", "description": "图片路径"},
                "left": {"type": "number", "description": "左边距（英寸）"},
                "top": {"type": "number", "description": "上边距（英寸）"},
                "width": {"type": "number", "description": "宽度（英寸，可选）"},
                "height": {"type": "number", "description": "高度（英寸，可选）"}
            },
            "required": ["session_id", "slide_index", "image_path", "left", "top"]
        }
    },
    {
        "name": "pptx_add_table",
        "description": "添加表格到幻灯片",
        "inputSchema": {
            "type": "object",
            "properties": {
                "session_id": {"type": "string"},
                "slide_index": {"type": "integer"},
                "rows": {"type": "integer", "description": "行数"},
                "cols": {"type": "integer", "description": "列数"},
                "left": {"type": "number", "default": 1},
                "top": {"type": "number", "default": 2},
                "width": {"type": "number", "default": 8},
                "height": {"type": "number", "default": 4},
                "data": {"type": "array", "description": "表格数据（可选）"}
            },
            "required": ["session_id", "slide_index", "rows", "cols"]
        }
    },
    {
        "name": "pptx_read_content",
        "description": "读取演示文稿中所有文本内容",
        "inputSchema": {
            "type": "object",
            "properties": {
                "session_id": {"type": "string"}
            },
            "required": ["session_id"]
        }
    },
    {
        "name": "pptx_list_slides",
        "description": "列出所有幻灯片信息",
        "inputSchema": {
            "type": "object",
            "properties": {
                "session_id": {"type": "string"}
            },
            "required": ["session_id"]
        }
    },
    {
        "name": "pptx_validate",
        "description": "验证 PPTX 文件安全性",
        "inputSchema": {
            "type": "object",
            "properties": {
                "file_path": {"type": "string", "description": "文件路径"}
            },
            "required": ["file_path"]
        }
    },
    {
        "name": "pptx_list_images",
        "description": "列出演示文稿中的所有图片，返回位置、尺寸、内容类型等信息",
        "inputSchema": {
            "type": "object",
            "properties": {
                "session_id": {"type": "string", "description": "会话 ID"},
                "slide_index": {"type": "integer", "minimum": 0, "description": "幻灯片索引（可选，不填则返回所有幻灯片的图片）"}
            },
            "required": ["session_id"]
        }
    },
    {
        "name": "pptx_export_images",
        "description": "导出演示文稿中的图片到临时目录，返回文件路径和元数据",
        "inputSchema": {
            "type": "object",
            "properties": {
                "session_id": {"type": "string", "description": "会话 ID"},
                "slide_index": {"type": "integer", "minimum": 0, "description": "幻灯片索引（可选）"}
            },
            "required": ["session_id"]
        }
    },
    {
        "name": "pptx_describe_slide",
        "description": "返回 slide 的结构化布局描述，含所有元素的位置、类型、文本、图片引用及布局分析",
        "inputSchema": {
            "type": "object",
            "properties": {
                "session_id": {"type": "string", "description": "会话 ID"},
                "slide_index": {"type": "integer", "minimum": 0, "description": "幻灯片索引（0-based）"}
            },
            "required": ["session_id", "slide_index"]
        }
    },
    {
        "name": "pptx_export_slide_snapshot",
        "description": "导出 slide 的结构化快照（含布局 JSON + 图片资源，不依赖 LibreOffice）",
        "inputSchema": {
            "type": "object",
            "properties": {
                "session_id": {"type": "string", "description": "会话 ID"},
                "slide_index": {"type": "integer", "minimum": 0, "description": "幻灯片索引（0-based）"}
            },
            "required": ["session_id", "slide_index"]
        }
    },
    {
        "name": "pptx_get_animation_info",
        "description": "获取 slide 的动画和 transition 信息，含动画顺序、触发方式、时长、延迟",
        "inputSchema": {
            "type": "object",
            "properties": {
                "session_id": {"type": "string", "description": "会话 ID"},
                "slide_index": {"type": "integer", "minimum": 0, "description": "幻灯片索引（0-based）"}
            },
            "required": ["session_id", "slide_index"]
        }
    }
]


# ===== 异常定义 =====
class PptxError(Exception):
    """PPTX 操作异常基类"""
    pass

class ValidationError(PptxError):
    """验证失败（用户错误）"""
    pass

class SecurityError(PptxError):
    """安全检查失败"""
    pass

class SessionError(PptxError):
    """会话错误"""
    pass


# ===== MCP Handler =====
class McpHandler(BaseHTTPRequestHandler):
    """MCP HTTP Handler"""
    
    # 类变量，由 main 设置
    session_manager: Optional[SessionManager] = None
    tools: Optional[PptxTools] = None
    event_loop: Optional[asyncio.AbstractEventLoop] = None
    token: Optional[str] = None
    
    def log_message(self, format, *args):
        log.info(format % args)
    
    def do_GET(self):
        """处理 GET 请求"""
        if self.path == "/health":
            self._handle_health()
        elif self.path == "/tools/list":
            self._send_json(None, result={"tools": TOOLS})
        elif self.path == "/stats":
            self._handle_stats()
        else:
            self.send_error(404)
    
    def do_POST(self):
        """处理 POST 请求"""
        if self.path != "/mcp":
            self.send_error(404)
            return
        
        # 认证检查
        if self.token:
            auth = self.headers.get("Authorization", "")
            if auth != f"Bearer {self.token}":
                self.send_error(401, "Unauthorized")
                return
        
        # 解析请求
        content_length = int(self.headers.get("Content-Length", 0))
        if content_length == 0:
            self.send_error(400, "Empty request")
            return
        
        body = self.rfile.read(content_length).decode("utf-8")
        
        try:
            request = json.loads(body)
        except json.JSONDecodeError:
            self.send_error(400, "Invalid JSON")
            return
        
        method = request.get("method", "")
        req_id = request.get("id")
        
        try:
            if method == "initialize":
                result = {
                    "protocolVersion": "2024-11-05",
                    "capabilities": {"tools": {}},
                    "serverInfo": {
                        "name": "pptx-server",
                        "version": "3.0.0"
                    }
                }
                self._send_json(req_id, result=result)
            elif method == "tools/list":
                self._send_json(req_id, result={"tools": TOOLS})
            elif method == "tools/call":
                result = self._handle_tool_call(request.get("params", {}))
                self._send_json(req_id, result=result)
            else:
                self._send_json(req_id, error={
                    "code": -32601,
                    "message": f"Unknown method: {method}"
                })
        except ValidationError as e:
            self._send_json(req_id, error={
                "code": -32001,
                "message": f"验证失败: {str(e)}"
            })
        except SecurityError as e:
            log.warning(f"Security violation: {e}")
            self._send_json(req_id, error={
                "code": -32002,
                "message": f"安全检查失败"
            })
        except SessionError as e:
            self._send_json(req_id, error={
                "code": -32003,
                "message": f"会话错误: {str(e)}"
            })
        except Exception as e:
            log.error(f"Internal error: {e}")
            self._send_json(req_id, error={
                "code": -32603,
                "message": f"内部错误: {str(e)}"
            })
    
    def _handle_health(self):
        """健康检查"""
        import shutil
        
        # 检查各项指标
        checks = {
            "session_manager": self.session_manager is not None,
            "tools": self.tools is not None,
            "temp_dir_writable": os.access(temp_manager.temp_dir, os.W_OK),
        }
        
        # 磁盘空间检查
        try:
            stat = shutil.disk_usage(temp_manager.temp_dir)
            checks["disk_space_ok"] = stat.free > 100 * 1024 * 1024
        except:
            checks["disk_space_ok"] = False
        
        is_healthy = all(checks.values())
        
        response = {
            "status": "healthy" if is_healthy else "degraded",
            "version": "3.0.0",
            "checks": checks,
            "limits": {
                "max_file_size_mb": limits.MAX_FILE_SIZE / 1024 / 1024,
                "max_slides": limits.MAX_SLIDES,
                "session_ttl_seconds": limits.SESSION_TTL
            }
        }
        
        if self.session_manager:
            response["stats"] = self.session_manager.get_stats()
        
        self._send_json(None, result=response)
    
    def _handle_stats(self):
        """统计信息"""
        response = {
            "sessions": self.session_manager.get_stats() if self.session_manager else {},
            "temp_files": temp_manager.get_stats()
        }
        self._send_json(None, result=response)
    
    def _handle_tool_call(self, params: dict) -> dict:
        """处理工具调用"""
        tool_name = params.get("name", "")
        args = params.get("arguments", {})
        
        # 工具路由
        handlers = {
            "pptx_create": self._tool_create,
            "pptx_open": self._tool_open,
            "pptx_save": self._tool_save,
            "pptx_close": self._tool_close,
            "pptx_info": self._tool_info,
            "pptx_add_slide": self._tool_add_slide,
            "pptx_add_text": self._tool_add_text,
            "pptx_add_image": self._tool_add_image,
            "pptx_add_table": self._tool_add_table,
            "pptx_read_content": self._tool_read_content,
            "pptx_list_slides": self._tool_list_slides,
            "pptx_validate": self._tool_validate,
            "pptx_list_images": self._tool_list_images,
            "pptx_export_images": self._tool_export_images,
            "pptx_describe_slide": self._tool_describe_slide,
            "pptx_export_slide_snapshot": self._tool_export_slide_snapshot,
            "pptx_get_animation_info": self._tool_get_animation_info,
        }
        
        if tool_name not in handlers:
            raise ValueError(f"未知工具: {tool_name}")
        
        # 执行工具（在事件循环中）
        result = handlers[tool_name](args)
        
        return {
            "content": [{
                "type": "text",
                "text": json.dumps(result, ensure_ascii=False, indent=2)
            }]
        }
    
    # ===== 工具实现（同步版本，线程安全）=====
    
    def _tool_create(self, args: dict) -> dict:
        return self.tools.create(
            name=args.get("name", "Untitled")
        )
    
    def _tool_open(self, args: dict) -> dict:
        return self.tools.open(
            file_path=args["file_path"]
        )
    
    def _tool_save(self, args: dict) -> dict:
        return self.tools.save(
            session_id=args["session_id"],
            output_path=args.get("output_path")
        )
    
    def _tool_close(self, args: dict) -> dict:
        return self.tools.close(
            session_id=args["session_id"]
        )
    
    def _tool_info(self, args: dict) -> dict:
        return self.tools.info(
            session_id=args["session_id"]
        )
    
    def _tool_add_slide(self, args: dict) -> dict:
        return self.tools.add_slide(
            session_id=args["session_id"],
            layout_index=args.get("layout_index", 0)
        )
    
    def _tool_add_text(self, args: dict) -> dict:
        return self.tools.add_text(
            session_id=args["session_id"],
            slide_index=args["slide_index"],
            text=args["text"],
            position=args.get("position", "body"),
            left=args.get("left", 1.0),
            top=args.get("top", 1.0),
            width=args.get("width", 8.0),
            height=args.get("height", 1.0),
            font_size=args.get("font_size", 18)
        )
    
    def _tool_add_image(self, args: dict) -> dict:
        return self.tools.add_image(
            session_id=args["session_id"],
            slide_index=args["slide_index"],
            image_path=args["image_path"],
            left=args["left"],
            top=args["top"],
            width=args.get("width"),
            height=args.get("height")
        )
    
    def _tool_add_table(self, args: dict) -> dict:
        return self.tools.add_table(
            session_id=args["session_id"],
            slide_index=args["slide_index"],
            rows=args["rows"],
            cols=args["cols"],
            left=args.get("left", 1.0),
            top=args.get("top", 2.0),
            width=args.get("width", 8.0),
            height=args.get("height", 4.0),
            data=args.get("data")
        )
    
    def _tool_read_content(self, args: dict) -> dict:
        return self.tools.read_content(
            session_id=args["session_id"]
        )
    
    def _tool_list_slides(self, args: dict) -> dict:
        return self.tools.list_slides(
            session_id=args["session_id"]
        )
    
    def _tool_validate(self, args: dict) -> dict:
        return self.tools.validate(
            file_path=args["file_path"]
        )

    def _tool_list_images(self, args: dict) -> dict:
        return self.tools.list_images(
            session_id=args["session_id"],
            slide_index=args.get("slide_index")
        )

    def _tool_export_images(self, args: dict) -> dict:
        return self.tools.export_images(
            session_id=args["session_id"],
            slide_index=args.get("slide_index")
        )

    def _tool_describe_slide(self, args: dict) -> dict:
        return self.tools.describe_slide(
            session_id=args["session_id"],
            slide_index=args["slide_index"]
        )

    def _tool_export_slide_snapshot(self, args: dict) -> dict:
        return self.tools.export_slide_snapshot(
            session_id=args["session_id"],
            slide_index=args["slide_index"]
        )

    def _tool_get_animation_info(self, args: dict) -> dict:
        return self.tools.get_animation_info(
            session_id=args["session_id"],
            slide_index=args["slide_index"]
        )

    def _send_json(self, req_id, result=None, error=None):
        """发送 JSON 响应"""
        response = {"jsonrpc": "2.0", "id": req_id}
        if error:
            response["error"] = error
        else:
            response["result"] = result
        
        body = json.dumps(response, ensure_ascii=False).encode("utf-8")
        self.send_response(200)
        self.send_header("Content-Type", "application/json")
        self.send_header("Content-Length", str(len(body)))
        self.end_headers()
        self.wfile.write(body)


class ThreadingHTTPServer(ThreadingMixIn, HTTPServer):
    """线程 HTTP 服务器"""
    daemon_threads = True


def main():
    """主函数"""
    # 解析参数
    parser = argparse.ArgumentParser(description="PPTX MCP Server")
    parser.add_argument("--port", type=int, default=8010, help="服务端口")
    parser.add_argument("--token", type=str, default="", help="认证令牌")
    parser.add_argument("--work-dir", type=str, default="/tmp/pptx-work", help="工作目录")
    args = parser.parse_args()
    
    # 创建工作目录
    work_dir = os.path.abspath(args.work_dir)
    os.makedirs(work_dir, exist_ok=True)
    
    # 设置类变量
    McpHandler.token = args.token if args.token else None
    
    log.info(f"Starting PPTX MCP Server v3.0")
    log.info(f"Port: {args.port}")
    log.info(f"Work directory: {work_dir}")
    log.info(f"Max file size: {limits.MAX_FILE_SIZE / 1024 / 1024}MB")
    log.info(f"Session TTL: {limits.SESSION_TTL}s")
    
    # 初始化会话管理器和工具
    session_manager = SessionManager()
    tools = PptxTools(session_manager, work_dir)
    
    McpHandler.session_manager = session_manager
    McpHandler.tools = tools
    
    # 启动会话管理器
    session_manager.start()
    
    # 启动 HTTP 服务器
    server = ThreadingHTTPServer(("0.0.0.0", args.port), McpHandler)
    
    log.info(f"Server listening on http://0.0.0.0:{args.port}")
    log.info(f"Health check: http://localhost:{args.port}/health")
    log.info(f"MCP endpoint: http://localhost:{args.port}/mcp")
    
    # 信号处理
    def shutdown(signum, frame):
        log.info("Shutting down...")
        session_manager.stop()
        server.shutdown()
        sys.exit(0)
    
    signal.signal(signal.SIGTERM, shutdown)
    signal.signal(signal.SIGINT, shutdown)
    
    try:
        server.serve_forever()
    except KeyboardInterrupt:
        pass
    finally:
        session_manager.stop()
        server.server_close()
        log.info("Server stopped")


if __name__ == "__main__":
    main()
