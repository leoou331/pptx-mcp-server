"""
会话管理模块
管理 PPTX 文档的打开/关闭/超时清理

修复：
- P0-1: 为每个 Session 添加独立锁
- P0-2: 使用 threading.Lock 替代 asyncio.Lock（线程安全）
"""
import os
import threading
import uuid
from datetime import datetime, timedelta
from typing import Optional, Dict, List
from dataclasses import dataclass, field
import logging

from .validator import validate_pptx, limits

log = logging.getLogger("pptx-server")


@dataclass
class Session:
    """会话数据（线程安全）"""
    id: str
    presentation: object  # Presentation 对象
    created_at: datetime
    last_accessed: datetime
    source_file: Optional[str] = None
    dirty: bool = False
    name: str = "Untitled"
    _lock: Optional[threading.Lock] = field(default=None, repr=False, init=False)
    
    def __post_init__(self):
        """初始化锁"""
        self._lock = threading.Lock()
    
    @property
    def lock(self) -> threading.Lock:
        """获取会话锁"""
        return self._lock


class SessionManager:
    """会话管理器（线程安全）"""
    
    def __init__(self):
        self._sessions: Dict[str, Session] = {}
        self._lock = threading.Lock()  # 线程安全锁
        self._cleanup_thread: Optional[threading.Thread] = None
        self._stop_event = threading.Event()
    
    def start(self):
        """启动清理线程"""
        self._cleanup_thread = threading.Thread(target=self._cleanup_loop, daemon=True)
        self._cleanup_thread.start()
        log.info("Session manager started")
    
    def stop(self):
        """停止清理线程"""
        self._stop_event.set()
        if self._cleanup_thread:
            self._cleanup_thread.join(timeout=5)
        log.info("Session manager stopped")
    
    def create(self, name: str = "Untitled") -> str:
        """
        创建新会话
        
        Args:
            name: 演示文稿名称
            
        Returns:
            会话 ID
        """
        from pptx import Presentation
        
        with self._lock:
            # 检查会话数量限制
            if len(self._sessions) >= 50:
                raise RuntimeError("活跃会话过多，请关闭不用的会话")
            
            session_id = str(uuid.uuid4())[:8]
            prs = Presentation()
            
            self._sessions[session_id] = Session(
                id=session_id,
                presentation=prs,
                created_at=datetime.now(),
                last_accessed=datetime.now(),
                name=name
            )
            
            log.info(f"Session created: {session_id} ({name})")
            return session_id
    
    def open(self, file_path: str, name: Optional[str] = None) -> str:
        """
        打开文件创建会话
        
        Args:
            file_path: 文件路径
            name: 演示文稿名称（可选）
            
        Returns:
            会话 ID
        """
        from pptx import Presentation
        
        # 安全验证
        valid, msg = validate_pptx(file_path)
        if not valid:
            raise ValueError(f"文件验证失败: {msg}")
        
        with self._lock:
            # 检查会话数量限制
            if len(self._sessions) >= 50:
                raise RuntimeError("活跃会话过多，请关闭不用的会话")
            
            session_id = str(uuid.uuid4())[:8]
            prs = Presentation(file_path)
            
            display_name = name or os.path.basename(file_path)
            
            self._sessions[session_id] = Session(
                id=session_id,
                presentation=prs,
                created_at=datetime.now(),
                last_accessed=datetime.now(),
                source_file=file_path,
                name=display_name
            )
            
            log.info(f"Session opened: {session_id} ({file_path})")
            return session_id
    
    def get(self, session_id: str) -> Session:
        """
        获取会话
        
        Args:
            session_id: 会话 ID
            
        Returns:
            Session 对象
            
        Raises:
            KeyError: 会话不存在
        """
        with self._lock:
            if session_id not in self._sessions:
                raise KeyError(f"会话不存在: {session_id}")
            
            session = self._sessions[session_id]
            session.last_accessed = datetime.now()
            return session
    
    def close(self, session_id: str) -> bool:
        """
        关闭会话
        
        Args:
            session_id: 会话 ID
            
        Returns:
            是否成功关闭
        """
        with self._lock:
            if session_id in self._sessions:
                del self._sessions[session_id]
                log.info(f"Session closed: {session_id}")
                return True
            return False
    
    def list_sessions(self) -> List[Dict]:
        """列出所有会话信息"""
        with self._lock:
            return [
                {
                    "id": s.id,
                    "name": s.name,
                    "slides": len(s.presentation.slides),
                    "dirty": s.dirty,
                    "created_at": s.created_at.isoformat(),
                    "last_accessed": s.last_accessed.isoformat(),
                    "source_file": s.source_file
                }
                for s in self._sessions.values()
            ]
    
    def _cleanup_loop(self):
        """定期清理过期会话"""
        while not self._stop_event.is_set():
            try:
                self._stop_event.wait(timeout=300)  # 每5分钟检查
                if self._stop_event.is_set():
                    break
                self._cleanup_expired()
            except Exception as e:
                log.error(f"Cleanup failed: {e}")
    
    def _cleanup_expired(self):
        """清理过期会话"""
        now = datetime.now()
        expired = []
        
        with self._lock:
            for sid, session in self._sessions.items():
                age = (now - session.last_accessed).total_seconds()
                if age > limits.SESSION_TTL:
                    expired.append(sid)
            
            for sid in expired:
                del self._sessions[sid]
                log.info(f"Session expired: {sid}")
        
        if expired:
            log.info(f"Cleaned up {len(expired)} expired sessions")
    
    def get_stats(self) -> dict:
        """获取统计信息"""
        with self._lock:
            return {
                "active_sessions": len(self._sessions),
                "max_sessions": 50,
                "session_ttl_seconds": limits.SESSION_TTL
            }
