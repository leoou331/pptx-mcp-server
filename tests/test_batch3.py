"""
Batch3 新增工具测试：
- pptx_manage_slide_masters
- pptx_apply_picture_effects
"""
import os
import pytest

from pptx import Presentation
from pptx.util import Inches
from pptx.oxml.ns import qn

from security.session import SessionManager
from tools.manager import PptxTools


@pytest.fixture
def tools_env(tmp_path):
    """创建工具实例和带幻灯片的会话。"""
    manager = SessionManager()
    tools = PptxTools(manager, str(tmp_path))
    session_id = tools.create("TestDeck")["session_id"]
    tools.add_slide(session_id, layout_index=0)
    return tools, session_id


@pytest.fixture
def image_env(tools_env, tmp_path):
    """创建带图片的幻灯片环境。"""
    tools, sid = tools_env

    # 创建一个小的测试图片
    from PIL import Image
    img_path = str(tmp_path / "test_image.png")
    img = Image.new("RGB", (100, 100), color="red")
    img.save(img_path)

    # 添加图片到幻灯片
    tools.add_image(
        session_id=sid,
        slide_index=0,
        image_path=img_path,
        left=1.0,
        top=1.0,
        width=3.0,
        height=2.0,
    )
    return tools, sid, img_path


# ===== pptx_manage_slide_masters =====

class TestManageSlideMasters:

    def test_list_masters(self, tools_env):
        tools, sid = tools_env
        result = tools.manage_slide_masters(session_id=sid, action="list")
        assert result["action"] == "list"
        assert result["total_masters"] >= 1
        assert "masters" in result
        # 每个 master 应该有 layouts 列表
        for master in result["masters"]:
            assert "index" in master
            assert "layout_count" in master
            assert "layouts" in master
            assert isinstance(master["layouts"], list)
            for layout in master["layouts"]:
                assert "index" in layout
                assert "name" in layout

    def test_list_masters_message(self, tools_env):
        tools, sid = tools_env
        result = tools.manage_slide_masters(session_id=sid, action="list")
        assert "找到" in result["message"]
        assert "母版" in result["message"]

    def test_apply_layout(self, tools_env):
        tools, sid = tools_env
        # 先获取可用的版式
        list_result = tools.manage_slide_masters(session_id=sid, action="list")
        masters = list_result["masters"]
        assert len(masters) > 0

        # 选择第一个 master 的最后一个 layout（不同于默认）
        layouts = masters[0]["layouts"]
        target_layout_idx = len(layouts) - 1 if len(layouts) > 1 else 0

        result = tools.manage_slide_masters(
            session_id=sid,
            action="apply",
            slide_index=0,
            master_index=0,
            layout_index=target_layout_idx,
        )
        assert result["action"] == "apply"
        assert result["slide_index"] == 0
        assert result["master_index"] == 0
        assert result["layout_index"] == target_layout_idx
        assert "layout_name" in result
        assert "已将幻灯片" in result["message"]

    def test_apply_invalid_action(self, tools_env):
        tools, sid = tools_env
        with pytest.raises(ValueError, match="无效的 action"):
            tools.manage_slide_masters(session_id=sid, action="delete")

    def test_apply_missing_slide_index(self, tools_env):
        tools, sid = tools_env
        with pytest.raises(ValueError, match="slide_index"):
            tools.manage_slide_masters(session_id=sid, action="apply", layout_index=0)

    def test_apply_missing_layout_index(self, tools_env):
        tools, sid = tools_env
        with pytest.raises(ValueError, match="layout_index"):
            tools.manage_slide_masters(session_id=sid, action="apply", slide_index=0)

    def test_apply_invalid_master_index(self, tools_env):
        tools, sid = tools_env
        with pytest.raises(ValueError, match="master_index.*超出范围"):
            tools.manage_slide_masters(
                session_id=sid,
                action="apply",
                slide_index=0,
                master_index=999,
                layout_index=0,
            )

    def test_apply_invalid_layout_index(self, tools_env):
        tools, sid = tools_env
        with pytest.raises(ValueError, match="layout_index.*超出范围"):
            tools.manage_slide_masters(
                session_id=sid,
                action="apply",
                slide_index=0,
                master_index=0,
                layout_index=999,
            )

    def test_apply_invalid_slide_index(self, tools_env):
        tools, sid = tools_env
        with pytest.raises(ValueError, match="slide_index.*超出范围"):
            tools.manage_slide_masters(
                session_id=sid,
                action="apply",
                slide_index=999,
                master_index=0,
                layout_index=0,
            )

    def test_negative_master_index(self, tools_env):
        tools, sid = tools_env
        with pytest.raises(ValueError, match="不能为负数"):
            tools.manage_slide_masters(
                session_id=sid, action="apply", master_index=-1,
                slide_index=0, layout_index=0,
            )

    def test_master_index_type_error(self, tools_env):
        tools, sid = tools_env
        with pytest.raises(TypeError, match="master_index 必须是整数"):
            tools.manage_slide_masters(
                session_id=sid, action="apply", master_index="abc",
                slide_index=0, layout_index=0,
            )

    def test_list_ignores_master_index(self, tools_env):
        """list 操作应忽略 master_index 参数，不进行验证。"""
        tools, sid = tools_env
        # 传入超大 master_index 也不应报错
        result = tools.manage_slide_masters(
            session_id=sid, action="list", master_index=999,
        )
        assert result["action"] == "list"
        assert result["total_masters"] >= 1

    def test_apply_removes_orphan_placeholders(self, tools_env):
        """切换到 Blank layout 后，孤立的 title/subtitle placeholder 应被移除。"""
        tools, sid = tools_env

        # 默认 layout_index=0 是 Title Slide，包含 title(idx=0) 和 subtitle(idx=1)
        # 先确认当前幻灯片有 placeholder
        session = tools.sessions.get(sid)
        slide = session.presentation.slides[0]
        initial_phs = [s for s in slide.shapes if s.is_placeholder]
        initial_ph_idxs = {
            s.placeholder_format.idx for s in slide.shapes if s.is_placeholder
        }
        assert 0 in initial_ph_idxs, "初始幻灯片应有 title placeholder (idx=0)"

        # 找到 Blank layout (index 6)，它没有 idx 0 和 1
        # 先确认 Blank layout 不包含 title placeholder
        master = session.presentation.slide_masters[0]
        blank_layout_idx = 6  # "Blank" layout in default theme
        blank_layout = list(master.slide_layouts)[blank_layout_idx]
        blank_ph_idxs = {
            ph.placeholder_format.idx for ph in blank_layout.placeholders
        }
        assert 0 not in blank_ph_idxs, "Blank layout should not have title placeholder"

        # 应用 Blank layout
        result = tools.manage_slide_masters(
            session_id=sid,
            action="apply",
            slide_index=0,
            master_index=0,
            layout_index=blank_layout_idx,
        )
        assert result["action"] == "apply"
        assert result["removed_placeholders"] > 0

        # 验证幻灯片上不再有 title/subtitle placeholder
        slide = session.presentation.slides[0]
        remaining_ph_idxs = {
            s.placeholder_format.idx for s in slide.shapes if s.is_placeholder
        }
        # idx 0 (title) 和 idx 1 (subtitle) 应被移除
        assert 0 not in remaining_ph_idxs, "Title placeholder (idx=0) should be removed"
        assert 1 not in remaining_ph_idxs, "Subtitle placeholder (idx=1) should be removed"
        # 共有的 placeholder (10, 11, 12) 应被保留
        for idx in blank_ph_idxs:
            if idx in initial_ph_idxs:
                assert idx in remaining_ph_idxs, (
                    f"Placeholder idx={idx} exists in both layouts and should be kept"
                )

    def test_apply_keeps_matching_placeholders(self, tools_env):
        """切换版式时，新版式中也有的 placeholder idx 应保留。"""
        tools, sid = tools_env

        session = tools.sessions.get(sid)
        master = session.presentation.slide_masters[0]

        # 找两个都有 placeholder 的版式（title slide -> title + content）
        layouts = list(master.slide_layouts)
        if len(layouts) < 2:
            pytest.skip("Need at least 2 layouts")

        # Layout 0 = Title Slide, Layout 1 = Title and Content（通常都有 title idx=0）
        # 先切到 layout 1
        result = tools.manage_slide_masters(
            session_id=sid,
            action="apply",
            slide_index=0,
            master_index=0,
            layout_index=1,
        )
        assert result["action"] == "apply"

        # 检查结果：layout 1 通常有 title (idx=0)，
        # 那么 title placeholder 应该被保留
        slide = session.presentation.slides[0]
        layout_1 = layouts[1]
        layout_1_ph_idxs = {
            ph.placeholder_format.idx for ph in layout_1.placeholders
        }
        remaining_ph_idxs = {
            s.placeholder_format.idx
            for s in slide.shapes if s.is_placeholder
        }
        # 所有保留的 placeholder 应该在新版式中存在
        assert remaining_ph_idxs.issubset(layout_1_ph_idxs), (
            f"Remaining placeholders {remaining_ph_idxs} should be subset "
            f"of new layout placeholders {layout_1_ph_idxs}"
        )


# ===== pptx_apply_picture_effects =====

class TestApplyPictureEffects:

    def test_crop(self, image_env):
        tools, sid, _ = image_env
        result = tools.apply_picture_effects(
            session_id=sid,
            slide_index=0,
            shape_index=0,
            effects={"crop": {"left": 0.1, "top": 0.2, "right": 0.1, "bottom": 0.2}},
        )
        assert result["applied_effects"] == ["crop"]
        assert result["effect_count"] == 1
        assert "已对图片" in result["message"]

    def test_border(self, image_env):
        tools, sid, _ = image_env
        result = tools.apply_picture_effects(
            session_id=sid,
            slide_index=0,
            shape_index=0,
            effects={"border": {"color": "FF0000", "width": 2.0}},
        )
        assert "border" in result["applied_effects"]

    def test_shadow_outer(self, image_env):
        tools, sid, _ = image_env
        result = tools.apply_picture_effects(
            session_id=sid,
            slide_index=0,
            shape_index=0,
            effects={"shadow": {
                "type": "outer",
                "blur_radius": 4.0,
                "distance": 3.0,
                "angle": 315,
                "color": "000000",
            }},
        )
        assert "shadow" in result["applied_effects"]

    def test_shadow_inner(self, image_env):
        tools, sid, _ = image_env
        result = tools.apply_picture_effects(
            session_id=sid,
            slide_index=0,
            shape_index=0,
            effects={"shadow": {"type": "inner", "blur_radius": 2.0}},
        )
        assert "shadow" in result["applied_effects"]

    def test_transparency(self, image_env):
        tools, sid, _ = image_env
        result = tools.apply_picture_effects(
            session_id=sid,
            slide_index=0,
            shape_index=0,
            effects={"transparency": 0.5},
        )
        assert "transparency" in result["applied_effects"]

    def test_brightness(self, image_env):
        tools, sid, _ = image_env
        result = tools.apply_picture_effects(
            session_id=sid,
            slide_index=0,
            shape_index=0,
            effects={"brightness": 0.3},
        )
        assert "brightness" in result["applied_effects"]

    def test_contrast(self, image_env):
        tools, sid, _ = image_env
        result = tools.apply_picture_effects(
            session_id=sid,
            slide_index=0,
            shape_index=0,
            effects={"contrast": -0.2},
        )
        assert "contrast" in result["applied_effects"]

    def test_multiple_effects(self, image_env):
        tools, sid, _ = image_env
        result = tools.apply_picture_effects(
            session_id=sid,
            slide_index=0,
            shape_index=0,
            effects={
                "crop": {"left": 0.05, "top": 0.05},
                "border": {"color": "0000FF", "width": 1.5},
                "shadow": {"type": "outer", "blur_radius": 3.0, "distance": 2.0},
                "transparency": 0.2,
                "brightness": 0.1,
                "contrast": 0.1,
            },
        )
        assert result["effect_count"] == 6
        assert set(result["applied_effects"]) == {
            "crop", "border", "shadow", "transparency", "brightness", "contrast",
        }

    def test_empty_effects_error(self, image_env):
        tools, sid, _ = image_env
        with pytest.raises(ValueError, match="effects 不能为空"):
            tools.apply_picture_effects(
                session_id=sid, slide_index=0, shape_index=0, effects={},
            )

    def test_invalid_effect_type(self, image_env):
        tools, sid, _ = image_env
        with pytest.raises(ValueError, match="不支持的效果类型"):
            tools.apply_picture_effects(
                session_id=sid, slide_index=0, shape_index=0,
                effects={"glow": {"size": 5}},
            )

    def test_crop_out_of_range(self, image_env):
        tools, sid, _ = image_env
        with pytest.raises(ValueError, match="crop.left 必须在 0.0~1.0"):
            tools.apply_picture_effects(
                session_id=sid, slide_index=0, shape_index=0,
                effects={"crop": {"left": 1.5}},
            )

    def test_transparency_out_of_range(self, image_env):
        tools, sid, _ = image_env
        with pytest.raises(ValueError, match="transparency 必须在 0.0~1.0"):
            tools.apply_picture_effects(
                session_id=sid, slide_index=0, shape_index=0,
                effects={"transparency": 2.0},
            )

    def test_brightness_out_of_range(self, image_env):
        tools, sid, _ = image_env
        with pytest.raises(ValueError, match="brightness 必须在 -1.0~1.0"):
            tools.apply_picture_effects(
                session_id=sid, slide_index=0, shape_index=0,
                effects={"brightness": -2.0},
            )

    def test_contrast_out_of_range(self, image_env):
        tools, sid, _ = image_env
        with pytest.raises(ValueError, match="contrast 必须在 -1.0~1.0"):
            tools.apply_picture_effects(
                session_id=sid, slide_index=0, shape_index=0,
                effects={"contrast": 1.5},
            )

    def test_no_pictures_on_slide(self, tools_env):
        tools, sid = tools_env
        with pytest.raises(ValueError, match="没有图片形状"):
            tools.apply_picture_effects(
                session_id=sid, slide_index=0, shape_index=0,
                effects={"crop": {"left": 0.1}},
            )

    def test_shape_index_out_of_range(self, image_env):
        tools, sid, _ = image_env
        with pytest.raises(ValueError, match="shape_index.*超出范围"):
            tools.apply_picture_effects(
                session_id=sid, slide_index=0, shape_index=99,
                effects={"crop": {"left": 0.1}},
            )

    def test_invalid_slide_index(self, image_env):
        tools, sid, _ = image_env
        with pytest.raises(ValueError, match="slide_index.*超出范围"):
            tools.apply_picture_effects(
                session_id=sid, slide_index=99, shape_index=0,
                effects={"crop": {"left": 0.1}},
            )

    def test_effects_not_dict(self, image_env):
        tools, sid, _ = image_env
        with pytest.raises(TypeError, match="effects 必须是字典"):
            tools.apply_picture_effects(
                session_id=sid, slide_index=0, shape_index=0, effects="invalid",
            )

    def test_negative_shape_index(self, image_env):
        tools, sid, _ = image_env
        with pytest.raises(ValueError, match="shape_index 不能为负数"):
            tools.apply_picture_effects(
                session_id=sid, slide_index=0, shape_index=-1,
                effects={"crop": {"left": 0.1}},
            )

    def test_invalid_shadow_type(self, image_env):
        tools, sid, _ = image_env
        with pytest.raises(ValueError, match="shadow.type 无效"):
            tools.apply_picture_effects(
                session_id=sid, slide_index=0, shape_index=0,
                effects={"shadow": {"type": "magic"}},
            )

    def test_border_invalid_color(self, image_env):
        tools, sid, _ = image_env
        with pytest.raises(ValueError, match="无效的颜色格式"):
            tools.apply_picture_effects(
                session_id=sid, slide_index=0, shape_index=0,
                effects={"border": {"color": "ZZZZZZ"}},
            )

    def test_border_negative_width(self, image_env):
        tools, sid, _ = image_env
        with pytest.raises(ValueError, match="border.width 必须大于 0"):
            tools.apply_picture_effects(
                session_id=sid, slide_index=0, shape_index=0,
                effects={"border": {"width": -1}},
            )

    def test_crop_partial_fields(self, image_env):
        """只设置部分裁剪字段（如仅 left 和 top）也应正常工作。"""
        tools, sid, _ = image_env
        result = tools.apply_picture_effects(
            session_id=sid,
            slide_index=0,
            shape_index=0,
            effects={"crop": {"left": 0.1, "top": 0.15}},
        )
        assert "crop" in result["applied_effects"]

    def test_shadow_defaults(self, image_env):
        """shadow 使用默认值（只传 type）也应正常工作。"""
        tools, sid, _ = image_env
        result = tools.apply_picture_effects(
            session_id=sid,
            slide_index=0,
            shape_index=0,
            effects={"shadow": {}},
        )
        assert "shadow" in result["applied_effects"]

    def test_brightness_and_contrast_combined(self, image_env):
        """同时设置亮度和对比度应正确合并到同一 lum 元素。"""
        tools, sid, _ = image_env
        result = tools.apply_picture_effects(
            session_id=sid,
            slide_index=0,
            shape_index=0,
            effects={"brightness": 0.2, "contrast": 0.3},
        )
        assert "brightness" in result["applied_effects"]
        assert "contrast" in result["applied_effects"]

        # 验证 XML：只有一个 a:lum 元素，同时包含 bright 和 contrast 属性
        session = tools.sessions.get(sid)
        prs = session.presentation
        slide = prs.slides[0]
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        pictures = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
        pic = pictures[0]
        bf = pic._element.find(qn("p:blipFill")); blipFill = bf if bf is not None else pic._element.find(qn("pic:blipFill"))
        blip = blipFill.find(qn("a:blip"))
        lum_elements = blip.findall(qn("a:lum"))
        assert len(lum_elements) == 1, f"期望 1 个 a:lum 元素，实际 {len(lum_elements)} 个"
        assert lum_elements[0].get("bright") == "20000"
        assert lum_elements[0].get("contrast") == "30000"

    def test_brightness_only_no_contrast_attr(self, image_env):
        """只设置 brightness 时，a:lum 不应包含 contrast 属性。"""
        tools, sid, _ = image_env
        tools.apply_picture_effects(
            session_id=sid, slide_index=0, shape_index=0,
            effects={"brightness": 0.5},
        )
        session = tools.sessions.get(sid)
        prs = session.presentation
        slide = prs.slides[0]
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        pictures = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
        bf = pictures[0]._element.find(qn("p:blipFill")); blipFill = bf if bf is not None else pictures[0]._element.find(qn("pic:blipFill"))
        blip = blipFill.find(qn("a:blip"))
        lum_elements = blip.findall(qn("a:lum"))
        assert len(lum_elements) == 1
        assert lum_elements[0].get("bright") == "50000"
        assert lum_elements[0].get("contrast") is None

    def test_contrast_zero_no_lum_element(self, image_env):
        """contrast == 0 且无 brightness 时，不应创建 a:lum 元素。"""
        tools, sid, _ = image_env
        tools.apply_picture_effects(
            session_id=sid, slide_index=0, shape_index=0,
            effects={"contrast": 0},
        )
        session = tools.sessions.get(sid)
        prs = session.presentation
        slide = prs.slides[0]
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        pictures = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
        bf = pictures[0]._element.find(qn("p:blipFill")); blipFill = bf if bf is not None else pictures[0]._element.find(qn("pic:blipFill"))
        blip = blipFill.find(qn("a:blip"))
        lum_elements = blip.findall(qn("a:lum"))
        assert len(lum_elements) == 0, "contrast=0 不应创建 a:lum 元素"

    def test_brightness_zero_contrast_nonzero(self, image_env):
        """brightness == 0, contrast != 0 时，只写 contrast 属性。"""
        tools, sid, _ = image_env
        tools.apply_picture_effects(
            session_id=sid, slide_index=0, shape_index=0,
            effects={"brightness": 0, "contrast": 0.4},
        )
        session = tools.sessions.get(sid)
        prs = session.presentation
        slide = prs.slides[0]
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        pictures = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
        bf = pictures[0]._element.find(qn("p:blipFill")); blipFill = bf if bf is not None else pictures[0]._element.find(qn("pic:blipFill"))
        blip = blipFill.find(qn("a:blip"))
        lum_elements = blip.findall(qn("a:lum"))
        assert len(lum_elements) == 1
        assert lum_elements[0].get("bright") is None
        assert lum_elements[0].get("contrast") == "40000"

    def test_perspective_shadow_rejected(self, image_env):
        """perspective 阴影类型应被拒绝。"""
        tools, sid, _ = image_env
        with pytest.raises(ValueError, match="shadow.type 无效.*perspective"):
            tools.apply_picture_effects(
                session_id=sid, slide_index=0, shape_index=0,
                effects={"shadow": {"type": "perspective"}},
            )

    def test_shadow_invalid_color_validated(self, image_env):
        """阴影颜色应经过验证。"""
        tools, sid, _ = image_env
        with pytest.raises(ValueError, match="无效的颜色格式"):
            tools.apply_picture_effects(
                session_id=sid, slide_index=0, shape_index=0,
                effects={"shadow": {"color": "ZZZZZZ"}},
            )
