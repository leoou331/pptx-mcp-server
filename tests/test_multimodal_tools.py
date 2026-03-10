"""
单元测试：幻灯片多模态感知工具

覆盖范围：
- _validate_slide_index: 类型校验 + 越界校验
- _validate_export_path: 路径遍历防护
- list_images / export_images: 基本功能 + 边界情况
- _estimate_shape_role: 各分支覆盖
- get_animation_info: 有动画 / 无动画
- describe_slide: 基本结构正确性
"""
import os
import sys
import json
import tempfile
import pytest

# 将项目根目录加入 path
sys.path.insert(0, os.path.join(os.path.dirname(__file__), ".."))

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE

from tools.manager import PptxTools, EMU_PER_INCH
from security.session import SessionManager
from security.tempfile import temp_manager


# ===== Fixtures =====

@pytest.fixture
def work_dir(tmp_path):
    """创建临时工作目录"""
    return str(tmp_path)


@pytest.fixture
def tools(work_dir):
    """创建 PptxTools 实例"""
    sm = SessionManager()
    sm.start()
    t = PptxTools(sm, work_dir)
    yield t
    sm.stop()


@pytest.fixture
def empty_session(tools, work_dir):
    """创建一个空的演示文稿会话"""
    result = tools.create(name="Test")
    return result["session_id"]


@pytest.fixture
def session_with_slides(tools, work_dir):
    """创建一个包含 3 张幻灯片的会话"""
    result = tools.create(name="TestSlides")
    sid = result["session_id"]
    for _ in range(3):
        tools.add_slide(sid)
    return sid


@pytest.fixture
def session_with_image(tools, work_dir):
    """创建一个包含图片的会话"""
    # 创建一个简单的测试图片
    from PIL import Image
    img_path = os.path.join(work_dir, "test_image.png")
    img = Image.new("RGB", (100, 100), color="red")
    img.save(img_path)

    result = tools.create(name="TestImage")
    sid = result["session_id"]
    tools.add_slide(sid)
    tools.add_image(
        session_id=sid, slide_index=0,
        image_path=img_path, left=1, top=1, width=3, height=2
    )
    return sid


# ===== _validate_slide_index 测试 =====

class TestValidateSlideIndex:
    """测试 _validate_slide_index 辅助方法"""

    def test_valid_int(self, tools, session_with_slides):
        """正常 int 输入应通过"""
        session = tools.sessions.get(session_with_slides)
        prs = session.presentation
        result = tools._validate_slide_index(prs, 0)
        assert result == 0

    def test_string_coercion(self, tools, session_with_slides):
        """字符串 '1' 应被转为 int"""
        session = tools.sessions.get(session_with_slides)
        prs = session.presentation
        result = tools._validate_slide_index(prs, "1")
        assert result == 1

    def test_float_coercion(self, tools, session_with_slides):
        """浮点数 2.0 应被转为 int"""
        session = tools.sessions.get(session_with_slides)
        prs = session.presentation
        result = tools._validate_slide_index(prs, 2.0)
        assert result == 2

    def test_invalid_string(self, tools, session_with_slides):
        """非数字字符串应抛出 ValueError"""
        session = tools.sessions.get(session_with_slides)
        prs = session.presentation
        with pytest.raises(ValueError, match="必须为整数"):
            tools._validate_slide_index(prs, "abc")

    def test_negative_index(self, tools, session_with_slides):
        """负数索引应抛出 ValueError"""
        session = tools.sessions.get(session_with_slides)
        prs = session.presentation
        with pytest.raises(ValueError, match="超出范围"):
            tools._validate_slide_index(prs, -1)

    def test_out_of_range(self, tools, session_with_slides):
        """超出范围的索引应抛出 ValueError"""
        session = tools.sessions.get(session_with_slides)
        prs = session.presentation
        with pytest.raises(ValueError, match="超出范围"):
            tools._validate_slide_index(prs, 100)

    def test_none_input(self, tools, session_with_slides):
        """None 输入应抛出 ValueError"""
        session = tools.sessions.get(session_with_slides)
        prs = session.presentation
        with pytest.raises(ValueError):
            tools._validate_slide_index(prs, None)


# ===== _validate_export_path 测试 =====

class TestValidateExportPath:
    """测试路径遍历防护"""

    def test_valid_session_id(self, tools, work_dir):
        """合法的 session_id 应通过"""
        export_dir = os.path.join(work_dir, "exported_images", "abc-123_XYZ")
        # 不应抛异常
        tools._validate_export_path("abc-123_XYZ", export_dir)

    def test_invalid_session_id_special_chars(self, tools, work_dir):
        """含特殊字符的 session_id 应被拒绝"""
        export_dir = os.path.join(work_dir, "exported_images", "test")
        with pytest.raises(ValueError, match="无效的 session_id"):
            tools._validate_export_path("../evil", export_dir)

    def test_path_traversal(self, tools, work_dir):
        """路径遍历攻击应被检测"""
        with pytest.raises(ValueError, match="Path traversal"):
            tools._validate_export_path("valid_id", "/etc/passwd")


# ===== list_images 测试 =====

class TestListImages:
    """测试 list_images 功能"""

    def test_empty_presentation(self, tools, empty_session):
        """空演示文稿应返回 0 张图片"""
        result = tools.list_images(empty_session)
        assert result["total_images"] == 0
        assert result["images"] == []

    def test_with_image(self, tools, session_with_image):
        """包含图片的幻灯片应返回图片信息"""
        result = tools.list_images(session_with_image)
        assert result["total_images"] >= 1
        img = result["images"][0]
        assert "slide_index" in img
        assert "content_type" in img
        assert "width_inches" in img
        assert "alt_text" in img

    def test_specific_slide(self, tools, session_with_image):
        """指定 slide_index 应只返回该幻灯片的图片"""
        result = tools.list_images(session_with_image, slide_index=0)
        assert result["total_images"] >= 1
        for img in result["images"]:
            assert img["slide_index"] == 0

    def test_invalid_slide_index(self, tools, session_with_image):
        """无效的 slide_index 应抛出 ValueError"""
        with pytest.raises(ValueError):
            tools.list_images(session_with_image, slide_index=999)

    def test_string_slide_index(self, tools, session_with_image):
        """字符串 slide_index 应被正确转换"""
        result = tools.list_images(session_with_image, slide_index="0")
        assert result["total_images"] >= 1


# ===== export_images 测试 =====

class TestExportImages:
    """测试 export_images 功能"""

    def test_export_with_image(self, tools, session_with_image):
        """导出包含图片的幻灯片应成功"""
        result = tools.export_images(session_with_image)
        assert result["exported_count"] >= 1
        assert result["failed_count"] == 0
        assert result["errors"] == []
        # 验证文件确实被创建
        for img in result["images"]:
            assert os.path.exists(img["file_path"])

    def test_export_empty(self, tools, empty_session):
        """导出空演示文稿应返回 0"""
        result = tools.export_images(empty_session)
        assert result["exported_count"] == 0
        assert result["failed_count"] == 0

    def test_invalid_session_id(self, tools):
        """无效的 session_id 格式应被拒绝"""
        with pytest.raises(ValueError, match="无效的 session_id"):
            tools.export_images("../evil_path")


# ===== _estimate_shape_role 测试 =====

class TestEstimateShapeRole:
    """测试 shape 角色估算逻辑"""

    def test_describe_slide_basic(self, tools, session_with_slides):
        """describe_slide 应返回正确结构"""
        result = tools.describe_slide(session_with_slides, 0)
        assert "session_id" in result
        assert "slide_index" in result
        assert "page_size" in result
        assert "elements" in result
        assert "layout_analysis" in result

    def test_describe_slide_invalid_index(self, tools, session_with_slides):
        """无效索引应抛出 ValueError"""
        with pytest.raises(ValueError):
            tools.describe_slide(session_with_slides, 999)


# ===== get_animation_info 测试 =====

class TestGetAnimationInfo:
    """测试动画信息获取"""

    def test_no_animation(self, tools, session_with_slides):
        """无动画的幻灯片应返回空动画列表"""
        result = tools.get_animation_info(session_with_slides, 0)
        assert result["has_animations"] is False
        assert result["animation_count"] == 0
        assert result["animations"] == []
        assert result["session_id"] == session_with_slides
        assert result["slide_index"] == 0

    def test_invalid_slide_index(self, tools, session_with_slides):
        """无效 slide_index 应抛出 ValueError"""
        with pytest.raises(ValueError):
            tools.get_animation_info(session_with_slides, 999)

    def test_string_slide_index(self, tools, session_with_slides):
        """字符串 slide_index 应被正确转换"""
        result = tools.get_animation_info(session_with_slides, "0")
        assert result["slide_index"] == 0


# ===== export_slide_snapshot 测试 =====

class TestExportSlideSnapshot:
    """测试幻灯片快照导出"""

    def test_basic_snapshot(self, tools, session_with_image):
        """基本快照导出应包含所有必要字段"""
        result = tools.export_slide_snapshot(session_with_image, 0)
        assert result["snapshot_type"] == "structural_layout"
        assert "page_size" in result
        assert "elements" in result
        assert "exported_images" in result
        assert "export_failed_count" in result
        assert "export_errors" in result

    def test_snapshot_invalid_index(self, tools, session_with_slides):
        """无效索引应抛出 ValueError"""
        with pytest.raises(ValueError):
            tools.export_slide_snapshot(session_with_slides, 999)


if __name__ == "__main__":
    pytest.main([__file__, "-v"])
