"""
PPTX 工具管理器
实现所有 MCP 工具（线程安全版本）

修复 P0-1: 使用 Session 级别的锁保护 Presentation 操作
"""
import os
import re
import logging
import tempfile
import warnings
from pathlib import Path
from typing import Dict, Any, List, Optional, Tuple

from PIL import Image, UnidentifiedImageError
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_SHAPE, PP_PLACEHOLDER
from pptx.enum.text import PP_ALIGN
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import ChartData
from pptx.oxml.ns import qn

from security.validator import validate_pptx, safe_path_in_dirs, limits
from security.session import SessionManager
from security.tempfile import temp_manager

log = logging.getLogger("pptx-server")

# 常量
EMU_PER_INCH = 914400


class PptxTools:
    """PPTX 工具集合（线程安全）"""

    MAX_NAME_LENGTH = 255
    MAX_SLIDE_COORDINATE_INCHES = 100.0
    MAX_FONT_SIZE = 400
    
    def __init__(self, session_manager: SessionManager, work_dir: str):
        self.sessions = session_manager
        self.work_dir = work_dir

    @staticmethod
    def _validate_slide_index(prs: Presentation, slide_index: int) -> int:
        """验证并规范化 slide_index 参数。

        将 slide_index 转为 int 并检查是否在有效范围 [0, len(slides)) 内。

        Args:
            prs: python-pptx Presentation 对象
            slide_index: 幻灯片索引（0-based）

        Returns:
            规范化后的 int 类型 slide_index

        Raises:
            ValueError: slide_index 类型无法转为 int 或超出范围
        """
        if not isinstance(slide_index, int):
            try:
                slide_index = int(slide_index)
            except (ValueError, TypeError):
                raise ValueError(
                    f"slide_index 必须为整数，收到 {type(slide_index).__name__}: {slide_index!r}"
                )
        total = len(prs.slides)
        if slide_index < 0 or slide_index >= total:
            raise ValueError(
                f"slide_index {slide_index} 超出范围 (有效范围: 0-{total - 1})"
            )
        return slide_index

    def _resolve_path(self, user_path: str) -> str:
        """将用户路径限制在工作目录或服务临时目录内。"""
        return safe_path_in_dirs(self.work_dir, user_path, temp_manager.temp_dir)

    def _validate_name(self, name: str) -> str:
        if not isinstance(name, str):
            raise TypeError("name 必须是字符串")

        cleaned = name.strip()
        if not cleaned:
            raise ValueError("name 不能为空")
        if len(cleaned) > self.MAX_NAME_LENGTH:
            raise ValueError(f"name 过长: {len(cleaned)} > {self.MAX_NAME_LENGTH}")
        return cleaned

    def _validate_non_negative(self, value: Any, field_name: str) -> float:
        if not isinstance(value, (int, float)) or isinstance(value, bool):
            raise TypeError(f"{field_name} 必须是数字")

        numeric = float(value)
        if numeric < 0:
            raise ValueError(f"{field_name} 不能为负数")
        if numeric > self.MAX_SLIDE_COORDINATE_INCHES:
            raise ValueError(
                f"{field_name} 超出允许范围: {numeric} > {self.MAX_SLIDE_COORDINATE_INCHES}"
            )
        return numeric

    def _validate_positive(self, value: Any, field_name: str) -> float:
        numeric = self._validate_non_negative(value, field_name)
        if numeric <= 0:
            raise ValueError(f"{field_name} 必须大于 0")
        return numeric

    def _validate_int(self, value: Any, field_name: str, *, minimum: int = 0) -> int:
        if not isinstance(value, int) or isinstance(value, bool):
            raise TypeError(f"{field_name} 必须是整数")
        if value < minimum:
            raise ValueError(f"{field_name} 不能小于 {minimum}")
        return value

    def _save_presentation_atomically(self, presentation: Presentation, output_path: str) -> None:
        """通过同目录临时文件 + 原子替换避免目标文件损坏。"""
        target = Path(output_path)
        target.parent.mkdir(parents=True, exist_ok=True)

        temp_fd, temp_name = tempfile.mkstemp(
            prefix=f".{target.stem}.",
            suffix=target.suffix or ".pptx",
            dir=str(target.parent),
        )
        os.close(temp_fd)

        try:
            presentation.save(temp_name)
            os.replace(temp_name, target)
        except Exception:
            try:
                os.unlink(temp_name)
            except FileNotFoundError:
                pass
            raise

    def _validate_image_file(self, image_path: str) -> None:
        """预先验证图片，避免将异常直接抛给 python-pptx/Pillow 内部栈。"""
        image_size = os.path.getsize(image_path)
        if image_size > limits.MAX_IMAGE_SIZE:
            raise ValueError(
                f"图片过大: {image_size/1024/1024:.1f}MB > {limits.MAX_IMAGE_SIZE/1024/1024}MB"
            )

        try:
            with warnings.catch_warnings():
                warnings.simplefilter("error", Image.DecompressionBombWarning)
                with Image.open(image_path) as image:
                    image.verify()
        except Image.DecompressionBombWarning as exc:
            raise ValueError(f"图片像素过大，疑似资源炸弹: {image_path}") from exc
        except (UnidentifiedImageError, OSError, ValueError) as exc:
            raise ValueError(f"无效图片文件: {image_path}") from exc
    
    def create(self, name: str = "Untitled") -> Dict[str, Any]:
        """
        创建空白演示文稿
        
        Args:
            name: 演示文稿名称
            
        Returns:
            会话信息
        """
        cleaned_name = self._validate_name(name)
        session_id = self.sessions.create(cleaned_name)
        
        return {
            "session_id": session_id,
            "name": cleaned_name,
            "slide_count": 0,
            "message": "演示文稿已创建"
        }
    
    def open(self, file_path: str) -> Dict[str, Any]:
        """
        打开现有文件
        
        Args:
            file_path: 文件路径
            
        Returns:
            会话信息
        """
        file_path = self._resolve_path(file_path)
        if not os.path.isfile(file_path):
            raise FileNotFoundError(f"文件不存在: {file_path}")
        
        session_id = self.sessions.open(file_path)
        session = self.sessions.get(session_id)
        
        with session.lock:  # 使用 Session 级别的锁
            return {
                "session_id": session_id,
                "name": session.name,
                "source_file": file_path,
                "slide_count": len(session.presentation.slides),
                "message": "文件已打开"
            }
    
    def save(self, session_id: str, output_path: Optional[str] = None) -> Dict[str, Any]:
        """
        保存演示文稿
        
        Args:
            session_id: 会话 ID
            output_path: 输出路径（可选）
            
        Returns:
            保存结果
        """
        session = self.sessions.get(session_id)
        
        if output_path:
            output_path = self._resolve_path(output_path)
        else:
            output_path = temp_manager.create(suffix='.pptx')
        
        with session.lock:  # 使用 Session 级别的锁
            prs = session.presentation
            self._save_presentation_atomically(prs, output_path)
            session.dirty = False
            
            return {
                "file_path": output_path,
                "slide_count": len(prs.slides),
                "message": "文件已保存"
            }
    
    def close(self, session_id: str) -> Dict[str, Any]:
        """
        关闭会话
        
        Args:
            session_id: 会话 ID
            
        Returns:
            关闭结果
        """
        closed = self.sessions.close(session_id)
        return {
            "closed": closed,
            "message": "会话已关闭" if closed else "会话不存在"
        }
    
    def info(self, session_id: str) -> Dict[str, Any]:
        """
        获取演示文稿信息
        
        Args:
            session_id: 会话 ID
            
        Returns:
            演示文稿信息
        """
        session = self.sessions.get(session_id)
        
        with session.lock:  # 使用 Session 级别的锁
            prs = session.presentation
            
            # 收集幻灯片信息
            slides_info = []
            for i, slide in enumerate(prs.slides):
                text_preview = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text_preview.append(shape.text[:100])
                
                slides_info.append({
                    "index": i,
                    "shape_count": len(slide.shapes),
                    "text_preview": " | ".join(text_preview)[:200]
                })
            
            return {
                "session_id": session_id,
                "name": session.name,
                "slide_count": len(prs.slides),
                "slide_width_inches": round(prs.slide_width.inches, 2),
                "slide_height_inches": round(prs.slide_height.inches, 2),
                "source_file": session.source_file,
                "dirty": session.dirty,
                "slides": slides_info
            }
    
    def add_slide(self, session_id: str, layout_index: int = 0) -> Dict[str, Any]:
        """
        添加幻灯片
        
        Args:
            session_id: 会话 ID
            layout_index: 布局索引
            
        Returns:
            添加结果
        """
        session = self.sessions.get(session_id)
        layout_index = self._validate_int(layout_index, "layout_index")
        
        with session.lock:  # 使用 Session 级别的锁
            prs = session.presentation
            
            # 检查幻灯片数量限制
            if len(prs.slides) >= limits.MAX_SLIDES:
                raise ValueError(f"幻灯片数量已达上限: {limits.MAX_SLIDES}")
            
            # 验证布局索引
            if layout_index >= len(prs.slide_layouts):
                raise ValueError(f"布局索引越界: {layout_index}")
            
            slide = prs.slides.add_slide(prs.slide_layouts[layout_index])
            session.dirty = True
            
            return {
                "slide_index": len(prs.slides) - 1,
                "layout_index": layout_index,
                "total_slides": len(prs.slides),
                "message": "幻灯片已添加"
            }
    
    def add_text(
        self,
        session_id: str,
        slide_index: int,
        text: str,
        position: str = "body",
        left: float = 1.0,
        top: float = 1.0,
        width: float = 8.0,
        height: float = 1.0,
        font_size: int = 18
    ) -> Dict[str, Any]:
        """
        添加文本
        
        Args:
            session_id: 会话 ID
            slide_index: 幻灯片索引
            text: 文本内容
            position: 位置类型 (title/body/custom)
            left, top, width, height: 位置和尺寸（英寸）
            font_size: 字号
            
        Returns:
            添加结果
        """
        session = self.sessions.get(session_id)
        slide_index = self._validate_int(slide_index, "slide_index")
        
        # 文本长度检查
        if not isinstance(text, str):
            raise TypeError("text 必须是字符串")
        if len(text) > limits.MAX_TEXT_LENGTH:
            raise ValueError(f"文本过长: {len(text)} > {limits.MAX_TEXT_LENGTH}")
        if position not in {"title", "body", "custom"}:
            raise ValueError(f"无效 position: {position}")

        left = self._validate_non_negative(left, "left")
        top = self._validate_non_negative(top, "top")
        width = self._validate_positive(width, "width")
        height = self._validate_positive(height, "height")
        font_size = self._validate_int(font_size, "font_size", minimum=1)
        if font_size > self.MAX_FONT_SIZE:
            raise ValueError(f"font_size 过大: {font_size} > {self.MAX_FONT_SIZE}")
        
        with session.lock:  # 使用 Session 级别的锁
            prs = session.presentation
            
            # 验证幻灯片索引
            if slide_index < 0 or slide_index >= len(prs.slides):
                raise ValueError(f"幻灯片索引越界: {slide_index}")
            
            slide = prs.slides[slide_index]
            
            # 根据位置类型处理
            if position == "title" and slide.shapes.title:
                slide.shapes.title.text = text
            else:
                # 创建文本框
                textbox = slide.shapes.add_textbox(
                    Inches(left),
                    Inches(top),
                    Inches(width),
                    Inches(height)
                )
                text_frame = textbox.text_frame
                text_frame.text = text
                
                # 设置字号
                if font_size:
                    for paragraph in text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.size = Pt(font_size)
            
            session.dirty = True
            
            return {
                "slide_index": slide_index,
                "position": position,
                "text_length": len(text),
                "message": "文本已添加"
            }
    
    def add_image(
        self,
        session_id: str,
        slide_index: int,
        image_path: str,
        left: float,
        top: float,
        width: Optional[float] = None,
        height: Optional[float] = None
    ) -> Dict[str, Any]:
        """
        添加图片
        
        Args:
            session_id: 会话 ID
            slide_index: 幻灯片索引
            image_path: 图片路径
            left, top: 位置（英寸）
            width, height: 尺寸（英寸，可选）
            
        Returns:
            添加结果
        """
        session = self.sessions.get(session_id)
        slide_index = self._validate_int(slide_index, "slide_index")
        left = self._validate_non_negative(left, "left")
        top = self._validate_non_negative(top, "top")
        if width is not None:
            width = self._validate_positive(width, "width")
        if height is not None:
            height = self._validate_positive(height, "height")
        
        image_path = self._resolve_path(image_path)
        
        # 检查文件存在
        if not os.path.isfile(image_path):
            raise FileNotFoundError(f"图片不存在: {image_path}")
        
        self._validate_image_file(image_path)
        image_size = os.path.getsize(image_path)
        
        with session.lock:  # 使用 Session 级别的锁
            prs = session.presentation
            
            # 验证幻灯片索引
            if slide_index < 0 or slide_index >= len(prs.slides):
                raise ValueError(f"幻灯片索引越界: {slide_index}")
            
            slide = prs.slides[slide_index]
            
            # 添加图片
            if width and height:
                slide.shapes.add_picture(
                    image_path,
                    Inches(left),
                    Inches(top),
                    Inches(width),
                    Inches(height)
                )
            elif width:
                slide.shapes.add_picture(
                    image_path,
                    Inches(left),
                    Inches(top),
                    width=Inches(width)
                )
            elif height:
                slide.shapes.add_picture(
                    image_path,
                    Inches(left),
                    Inches(top),
                    height=Inches(height)
                )
            else:
                slide.shapes.add_picture(
                    image_path,
                    Inches(left),
                    Inches(top)
                )
            
            session.dirty = True
            
            return {
                "slide_index": slide_index,
                "image_path": image_path,
                "image_size_kb": round(image_size / 1024, 1),
                "message": "图片已添加"
            }
    
    def add_table(
        self,
        session_id: str,
        slide_index: int,
        rows: int,
        cols: int,
        left: float = 1.0,
        top: float = 2.0,
        width: float = 8.0,
        height: float = 4.0,
        data: Optional[list] = None
    ) -> Dict[str, Any]:
        """
        添加表格
        
        Args:
            session_id: 会话 ID
            slide_index: 幻灯片索引
            rows, cols: 行数和列数
            left, top, width, height: 位置和尺寸（英寸）
            data: 表格数据（可选）
            
        Returns:
            添加结果
        """
        session = self.sessions.get(session_id)
        slide_index = self._validate_int(slide_index, "slide_index")
        rows = self._validate_int(rows, "rows", minimum=1)
        cols = self._validate_int(cols, "cols", minimum=1)
        left = self._validate_non_negative(left, "left")
        top = self._validate_non_negative(top, "top")
        width = self._validate_positive(width, "width")
        height = self._validate_positive(height, "height")

        if rows > limits.MAX_TABLE_ROWS:
            raise ValueError(f"rows 过大: {rows} > {limits.MAX_TABLE_ROWS}")
        if cols > limits.MAX_TABLE_COLS:
            raise ValueError(f"cols 过大: {cols} > {limits.MAX_TABLE_COLS}")
        if rows * cols > limits.MAX_TABLE_CELLS:
            raise ValueError(f"表格单元格过多: {rows * cols} > {limits.MAX_TABLE_CELLS}")
        if data is not None and not isinstance(data, list):
            raise TypeError("data 必须是二维数组")
        
        with session.lock:  # 使用 Session 级别的锁
            prs = session.presentation
            
            # 验证幻灯片索引
            if slide_index < 0 or slide_index >= len(prs.slides):
                raise ValueError(f"幻灯片索引越界: {slide_index}")
            
            slide = prs.slides[slide_index]
            
            # 添加表格
            table = slide.shapes.add_table(
                rows, cols,
                Inches(left),
                Inches(top),
                Inches(width),
                Inches(height)
            ).table
            
            # 填充数据
            if data:
                for i, row_data in enumerate(data[:rows]):
                    if row_data is None:
                        continue
                    if not isinstance(row_data, list):
                        raise TypeError("data 的每一行都必须是数组")
                    for j, cell_data in enumerate(row_data[:cols]):
                        if cell_data is not None:
                            table.cell(i, j).text = str(cell_data)
            
            session.dirty = True
            
            return {
                "slide_index": slide_index,
                "rows": rows,
                "cols": cols,
                "has_data": data is not None,
                "message": "表格已添加"
            }
    
    def read_content(self, session_id: str) -> Dict[str, Any]:
        """
        读取所有文本内容
        
        Args:
            session_id: 会话 ID
            
        Returns:
            所有幻灯片的文本内容
        """
        session = self.sessions.get(session_id)
        
        with session.lock:  # 使用 Session 级别的锁
            prs = session.presentation
            
            content = []
            for i, slide in enumerate(prs.slides):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text.append(shape.text)
                
                content.append({
                    "slide_index": i,
                    "text": "\n".join(slide_text),
                    "char_count": sum(len(t) for t in slide_text)
                })
            
            return {
                "session_id": session_id,
                "total_slides": len(prs.slides),
                "total_chars": sum(c["char_count"] for c in content),
                "slides": content
            }
    
    def list_slides(self, session_id: str) -> Dict[str, Any]:
        """
        列出所有幻灯片
        
        Args:
            session_id: 会话 ID
            
        Returns:
            幻灯片列表
        """
        session = self.sessions.get(session_id)
        
        with session.lock:  # 使用 Session 级别的锁
            prs = session.presentation
            
            slides = []
            for i, slide in enumerate(prs.slides):
                shapes_info = []
                for shape in slide.shapes:
                    shape_type = str(shape.shape_type)
                    has_text = hasattr(shape, "text") and bool(shape.text)
                    
                    shape_info = {
                        "type": shape_type,
                        "has_text": has_text,
                        "text_preview": shape.text[:50] if has_text and shape.text else None
                    }
                    shapes_info.append(shape_info)
                
                slides.append({
                    "index": i,
                    "shape_count": len(slide.shapes),
                    "shapes": shapes_info[:10]  # 限制返回数量
                })
            
            return {
                "session_id": session_id,
                "total_slides": len(prs.slides),
                "slides": slides
            }
    
    def validate(self, file_path: str) -> Dict[str, Any]:
        """
        验证文件安全性
        
        Args:
            file_path: 文件路径
            
        Returns:
            验证结果
        """
        file_path = self._resolve_path(file_path)
        if not os.path.isfile(file_path):
            raise FileNotFoundError(f"文件不存在: {file_path}")
        
        valid, message = validate_pptx(file_path)
        
        result = {
            "file_path": file_path,
            "valid": valid,
            "message": message
        }
        
        if valid:
            # 添加文件信息
            try:
                prs = Presentation(file_path)
                result["slide_count"] = len(prs.slides)
                result["slide_width_inches"] = round(prs.slide_width.inches, 2)
                result["slide_height_inches"] = round(prs.slide_height.inches, 2)
            except Exception as e:
                result["parse_error"] = str(e)
        
        return result


    def _validate_export_path(self, session_id: str, export_dir: str) -> None:
        """验证导出路径安全性，防止路径遍历攻击。

        Args:
            session_id: 会话 ID（只允许字母数字和 _-）
            export_dir: 导出目录路径

        Raises:
            ValueError: session_id 格式无效或路径遍历检测失败
        """
        if not re.match(r"^[a-zA-Z0-9_-]+$", session_id):
            raise ValueError(f"无效的 session_id: {session_id}")
        real_ed = os.path.realpath(export_dir)
        real_wd = os.path.realpath(self.work_dir)
        if not real_ed.startswith(real_wd + os.sep):
            raise ValueError("Path traversal detected")

    def _iter_picture_shapes(self, shapes) -> list:
        """递归遍历 shapes 集合，yield 所有包含图片的 shape（含 Group 内嵌套）。

        Args:
            shapes: python-pptx shapes 集合

        Yields:
            包含图片数据的 shape 对象
        """
        for shape in shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                yield shape
            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                yield from self._iter_picture_shapes(shape.shapes)
            else:
                if hasattr(shape, "image"):
                    try:
                        _ = shape.image
                        yield shape
                    except Exception:
                        pass

    def _build_image_info(self, s_idx: int, sh_idx: int, shape, nested_index: Optional[int] = None) -> Dict[str, Any]:
        """构建单个图片 shape 的元数据信息。

        Args:
            s_idx: 幻灯片索引
            sh_idx: shape 在幻灯片中的索引（group shape 的索引）
            shape: python-pptx shape 对象
            nested_index: group 内嵌套图片的子索引（可选），用于区分同一 group 中的不同图片

        Returns:
            包含图片位置、尺寸、alt text 等信息的字典
        """
        try:
            ct = shape.image.content_type
        except Exception:
            ct = "unknown"
        alt = shape.name or ""
        try:
            # 搜索 PML 和 DrawingML 命名空间下的 cNvPr 以获取 alt text
            el = shape._element.find(f".//{qn('p:cNvPr')}")
            if el is None:
                el = shape._element.find(f".//{qn('a:cNvPr')}")
            if el is not None:
                alt = el.get("descr", shape.name or "")
        except Exception:
            pass
        # Use a composite index to uniquely identify nested pictures inside
        # group shapes.  Format: "<group_index>.<nested_index>" when the
        # picture comes from a group, otherwise the plain shape index.
        if nested_index is not None:
            unique_index = f"{sh_idx}.{nested_index}"
        else:
            unique_index = sh_idx
        return {
            "slide_index": s_idx, "shape_index": unique_index, "name": shape.name,
            "content_type": ct,
            "left_inches": round(shape.left / EMU_PER_INCH, 4) if shape.left else 0,
            "top_inches": round(shape.top / EMU_PER_INCH, 4) if shape.top else 0,
            "width_inches": round(shape.width / EMU_PER_INCH, 4) if shape.width else 0,
            "height_inches": round(shape.height / EMU_PER_INCH, 4) if shape.height else 0,
            "z_order": sh_idx, "alt_text": alt,
        }

    def list_images(self, session_id: str, slide_index: Optional[int] = None) -> Dict[str, Any]:
        """列出演示文稿中的所有图片，返回位置、尺寸、内容类型等信息。

        Args:
            session_id: 会话 ID
            slide_index: 幻灯片索引（可选，不填则返回所有幻灯片的图片）

        Returns:
            包含图片列表和总数的字典
        """
        session = self.sessions.get(session_id)
        with session.lock:
            prs = session.presentation
            if slide_index is not None:
                slide_index = self._validate_slide_index(prs, slide_index)
                slides_to_check = [(slide_index, prs.slides[slide_index])]
            else:
                slides_to_check = list(enumerate(prs.slides))
            images = []
            for s_idx, slide in slides_to_check:
                for sh_idx, shape in enumerate(slide.shapes):
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        images.append(self._build_image_info(s_idx, sh_idx, shape))
                    elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                        for ni, nested in enumerate(self._iter_picture_shapes(shape.shapes)):
                            images.append(self._build_image_info(s_idx, sh_idx, nested, nested_index=ni))
                    elif hasattr(shape, "image"):
                        try:
                            _ = shape.image
                            images.append(self._build_image_info(s_idx, sh_idx, shape))
                        except Exception:
                            pass
            return {"session_id": session_id, "total_images": len(images), "images": images}

    def _export_inner(self, prs: Presentation, session_id: str,
                       slide_index: Optional[int] = None) -> Tuple[List[Dict[str, Any]], List[Dict[str, Any]]]:
        """内部方法：将图片 blob 导出到工作目录。

        Args:
            prs: python-pptx Presentation 对象
            session_id: 会话 ID
            slide_index: 幻灯片索引（可选）

        Returns:
            (exported, errors) 元组：
            - exported: 成功导出的图片信息列表
            - errors: 导出失败的错误信息列表
        """
        ext_map = {
            "image/jpeg": "jpg", "image/jpg": "jpg", "image/png": "png",
            "image/gif": "gif", "image/bmp": "bmp", "image/tiff": "tiff",
            "image/x-emf": "emf", "image/x-wmf": "wmf",
        }
        if slide_index is not None:
            slide_index = self._validate_slide_index(prs, slide_index)
            its = [(slide_index, prs.slides[slide_index])]
        else:
            its = list(enumerate(prs.slides))
        export_dir = os.path.join(self.work_dir, "exported_images", session_id)
        self._validate_export_path(session_id, export_dir)
        os.makedirs(export_dir, exist_ok=True)
        out: List[Dict[str, Any]] = []
        errors: List[Dict[str, Any]] = []
        for s_idx, slide in its:
            for sh_idx, shape in enumerate(slide.shapes):
                pics: List[Tuple[str, Any]] = []
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    pics.append((sh_idx, shape))
                elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    for ni, n in enumerate(self._iter_picture_shapes(shape.shapes)):
                        pics.append((f"{sh_idx}.{ni}", n))
                elif hasattr(shape, "image"):
                    try:
                        _ = shape.image
                        pics.append((sh_idx, shape))
                    except Exception:
                        pass
                for oi, ps in pics:
                    try:
                        blob = ps.image.blob
                        ct = ps.image.content_type
                        ext = ext_map.get(ct.lower(), "bin")
                        raw_name = f"slide{s_idx}_shape{oi}_{ps.name}.{ext}"
                        fn = "".join(
                            c if c.isalnum() or c in "._-" else "_"
                            for c in raw_name
                        )[:200]  # 限制文件名长度
                        fp = os.path.join(export_dir, fn)
                        with open(fp, "wb") as f:
                            f.write(blob)
                        # 注册到 temp_manager 以便自动清理
                        try:
                            temp_manager.register(fp)
                        except Exception:
                            pass
                        out.append({
                            "slide_index": s_idx, "shape_index": oi, "name": ps.name,
                            "content_type": ct, "file_path": fp,
                            "left_inches": round(ps.left / EMU_PER_INCH, 4) if ps.left else 0,
                            "top_inches": round(ps.top / EMU_PER_INCH, 4) if ps.top else 0,
                            "width_inches": round(ps.width / EMU_PER_INCH, 4) if ps.width else 0,
                            "height_inches": round(ps.height / EMU_PER_INCH, 4) if ps.height else 0,
                        })
                    except Exception as e:
                        log.warning(f"导出图片失败 slide={s_idx} shape={oi}: {e}")
                        errors.append({
                            "slide_index": s_idx,
                            "shape_index": oi,
                            "name": getattr(ps, "name", "unknown"),
                            "error": str(e),
                        })
        return out, errors

    def export_images(self, session_id: str, slide_index: Optional[int] = None) -> Dict[str, Any]:
        """导出图片到工作目录，通过 temp_manager 管理文件生命周期。

        Args:
            session_id: 会话 ID
            slide_index: 幻灯片索引（可选）

        Returns:
            包含导出结果、失败数量和错误详情的字典
        """
        if not re.match(r"^[a-zA-Z0-9_-]+$", session_id):
            raise ValueError(f"无效的 session_id: {session_id}")
        session = self.sessions.get(session_id)
        with session.lock:
            exported, errors = self._export_inner(session.presentation, session_id, slide_index)
        return {
            "session_id": session_id,
            "exported_count": len(exported),
            "failed_count": len(errors),
            "images": exported,
            "errors": errors,
        }

    def _estimate_shape_role(self, shape, bbox: Dict[str, float],
                             pw: float, ph: float) -> str:
        """根据 shape 类型、位置和尺寸估计其语义角色。

        Args:
            shape: python-pptx shape 对象
            bbox: 包含 left/top/width/height 的边界框字典（单位：英寸）
            pw: 页面宽度（英寸）
            ph: 页面高度（英寸）

        Returns:
            估计的角色字符串，如 'title', 'body', 'image' 等
        """
        st = shape.shape_type
        if st == MSO_SHAPE_TYPE.TABLE:
            return "table"
        if st == MSO_SHAPE_TYPE.CHART:
            return "chart"
        if st == MSO_SHAPE_TYPE.GROUP:
            return "group"
        if st == MSO_SHAPE_TYPE.PICTURE:
            ar = (bbox["width"] * bbox["height"]) / max(pw * ph, 0.001)
            return "hero_image" if ar > 0.3 else ("icon_or_logo" if ar < 0.05 else "image")
        if hasattr(shape, "text") and shape.text:
            text = shape.text.strip()
            try:
                if hasattr(shape, "placeholder_format") and shape.placeholder_format:
                    pt = shape.placeholder_format.type
                    if pt in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
                        return "title"
                    if pt == PP_PLACEHOLDER.SUBTITLE:
                        return "subtitle"
                    if pt == PP_PLACEHOLDER.BODY:
                        return "body"
            except Exception:
                pass
            if bbox["top"] < ph * 0.2 and bbox["height"] < ph * 0.15:
                return "title"
            if bbox["top"] < ph * 0.35 and len(text) < 100:
                return "subtitle_or_heading"
            ar = (bbox["width"] * bbox["height"]) / max(pw * ph, 0.001)
            return "body" if ar > 0.2 else "caption_or_label"
        return "decorative_shape"

    def _analyze_layout(self, elements: List[Dict[str, Any]],
                        pw: float, ph: float) -> Dict[str, Any]:
        """分析幻灯片元素的布局特征，包括阅读顺序、密度和重叠。

        注意：density_score 和 whitespace_ratio 是近似值，不考虑元素间的重叠。
        当多个元素重叠时，实际覆盖面积小于各元素面积之和，因此密度可能被高估。

        Args:
            elements: 元素信息列表（每个元素需包含 'bbox' 字段）
            pw: 页面宽度（英寸）
            ph: 页面高度（英寸）

        Returns:
            包含 reading_order, whitespace_ratio, density_score, overlaps 的字典
        """
        if not elements:
            return {"reading_order": [], "whitespace_ratio": 1.0, "density_score": 0.0, "overlaps": []}
        def rk(i: int) -> Tuple[int, float]:
            b = elements[i]["bbox"]
            return (int(b["top"] / max(ph / 10, 0.001)), b["left"] / max(pw, 0.001))
        ro = sorted(range(len(elements)), key=rk)
        # 近似计算：面积之和，未扣除重叠区域
        cov = sum(e["bbox"]["width"] * e["bbox"]["height"] for e in elements)
        d = min(1.0, cov / max(pw * ph, 0.001))
        def ov(b1: Dict, b2: Dict) -> bool:
            return not (
                b1["left"] + b1["width"] <= b2["left"] or
                b2["left"] + b2["width"] <= b1["left"] or
                b1["top"] + b1["height"] <= b2["top"] or
                b2["top"] + b2["height"] <= b1["top"]
            )
        # 限制重叠检测数量防止 O(n²) 性能问题
        MAX_OVERLAP_ELEMENTS = 50
        check_elems = elements[:MAX_OVERLAP_ELEMENTS]
        ovs = [
            [i, j]
            for i in range(len(check_elems))
            for j in range(i + 1, len(check_elems))
            if ov(check_elems[i]["bbox"], check_elems[j]["bbox"])
        ]
        ovs = ovs[:100]  # 最多返回 100 对
        return {
            "reading_order": ro,
            "whitespace_ratio": round(max(0, 1 - d), 3),
            "density_score": round(d, 3),
            "overlaps": ovs,
        }

    def _describe_inner(self, prs: Presentation, slide_index: int) -> Dict[str, Any]:
        """内部方法：构建幻灯片的结构化布局描述。

        Args:
            prs: python-pptx Presentation 对象
            slide_index: 幻灯片索引

        Returns:
            包含元素信息、布局分析等的结构化字典
        """
        slide_index = self._validate_slide_index(prs, slide_index)
        slide = prs.slides[slide_index]
        pw = (prs.slide_width or 9144000) / EMU_PER_INCH
        ph = (prs.slide_height or 6858000) / EMU_PER_INCH
        bg = {"type": "default", "color": None}
        try:
            fill = slide.background.fill
            if fill.type is not None:
                bg["type"] = str(fill.type)
                try:
                    bg["color"] = f"#{fill.fore_color.rgb}"
                except Exception:
                    pass
        except Exception:
            pass
        elements = []
        for sh_idx, shape in enumerate(slide.shapes):
            bbox = {
                "left": round(shape.left / EMU_PER_INCH, 4) if shape.left else 0,
                "top": round(shape.top / EMU_PER_INCH, 4) if shape.top else 0,
                "width": round(shape.width / EMU_PER_INCH, 4) if shape.width else 0,
                "height": round(shape.height / EMU_PER_INCH, 4) if shape.height else 0,
            }
            text_content = ""
            font_info = None
            if hasattr(shape, "text") and shape.text:
                text_content = shape.text[:500]
                try:
                    best_run, best_sz = None, -1
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            sz = run.font.size.pt if run.font.size else 0
                            if sz > best_sz:
                                best_sz, best_run = sz, run
                    if best_run:
                        font_info = {
                            "size_pt": best_run.font.size.pt if best_run.font.size else None,
                            "bold": best_run.font.bold,
                            "italic": best_run.font.italic,
                        }
                except Exception:
                    pass
            img_ref = None
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    img_ref = {
                        "content_type": shape.image.content_type,
                        "size_bytes": len(shape.image.blob),
                    }
                except Exception:
                    img_ref = {"content_type": "unknown", "size_bytes": 0}
            elements.append({
                "shape_index": sh_idx, "type": str(shape.shape_type), "name": shape.name,
                "bbox": bbox, "z_order": sh_idx, "text": text_content,
                "font_info": font_info, "image_ref": img_ref,
                "estimated_role": self._estimate_shape_role(shape, bbox, pw, ph),
            })
        return {
            "slide_index": slide_index,
            "page_size": {"width_inches": round(pw, 4), "height_inches": round(ph, 4)},
            "background": bg,
            "element_count": len(elements),
            "elements": elements,
            "layout_analysis": self._analyze_layout(elements, pw, ph),
        }

    def describe_slide(self, session_id: str, slide_index: int) -> Dict[str, Any]:
        """返回 slide 的结构化布局描述，含所有元素的位置、类型、文本、图片引用及布局分析。

        Args:
            session_id: 会话 ID
            slide_index: 幻灯片索引（0-based）

        Returns:
            结构化布局描述字典
        """
        session = self.sessions.get(session_id)
        with session.lock:
            slide_index = self._validate_slide_index(session.presentation, slide_index)
            r = self._describe_inner(session.presentation, slide_index)
            r["session_id"] = session_id
            return r

    def export_slide_snapshot(self, session_id: str, slide_index: int) -> Dict[str, Any]:
        """导出 slide 结构化快照（单次加锁，fallback 方案）。

        包含布局 JSON + 图片资源导出，不依赖 LibreOffice。

        Args:
            session_id: 会话 ID
            slide_index: 幻灯片索引（0-based）

        Returns:
            包含布局描述、导出图片和错误信息的字典
        """
        if not re.match(r"^[a-zA-Z0-9_-]+$", session_id):
            raise ValueError(f"无效的 session_id: {session_id}")
        session = self.sessions.get(session_id)
        with session.lock:
            prs = session.presentation
            slide_index = self._validate_slide_index(prs, slide_index)
            desc = self._describe_inner(prs, slide_index)
            imgs, errors = self._export_inner(prs, session_id, slide_index)
        return {
            "session_id": session_id, "slide_index": slide_index,
            "snapshot_type": "structural_layout",
            "note": "PNG rendering requires LibreOffice; returning structural layout JSON + image exports as fallback",
            "page_size": desc["page_size"], "background": desc["background"],
            "element_count": desc["element_count"], "elements": desc["elements"],
            "layout_analysis": desc["layout_analysis"],
            "exported_images": imgs,
            "export_failed_count": len(errors),
            "export_errors": errors,
        }

    def get_animation_info(self, session_id: str, slide_index: int) -> Dict[str, Any]:
        """获取 slide 动画和 transition 信息（通过解析 PML XML）。

        Args:
            session_id: 会话 ID
            slide_index: 幻灯片索引（0-based）

        Returns:
            包含动画列表、transition 信息等的字典
        """
        # nodeType 到 trigger 的映射表（基于 OOXML 规范）
        NODE_TYPE_TO_TRIGGER = {
            "clickEffect": "onClick",
            "withEffect": "withPrevious",
            "afterEffect": "afterPrevious",
            "mainSeq": "mainSequence",
            "interactiveSeq": "interactive",
        }
        session = self.sessions.get(session_id)
        with session.lock:
            prs = session.presentation
            slide_index = self._validate_slide_index(prs, slide_index)
            slide = prs.slides[slide_index]
            se = slide._element
            PML = "http://schemas.openxmlformats.org/presentationml/2006/main"
            te = se.find(f"{{{PML}}}transition")
            has_t = te is not None
            ti = None
            if has_t:
                ti = {
                    "type": "unknown",
                    "duration_ms": None,
                    "advance_on_click": te.get("advClick", "true").lower() != "false",
                    "advance_after_time_ms": None,
                }
                dur = te.get("dur")
                if dur:
                    try:
                        ti["duration_ms"] = int(dur)
                    except ValueError:
                        ti["duration_ms"] = dur
                adv = te.get("advTm")
                if adv:
                    try:
                        ti["advance_after_time_ms"] = int(adv)
                    except ValueError:
                        ti["advance_after_time_ms"] = adv
                for ch in te:
                    t = ch.tag.split("}")[-1] if "}" in ch.tag else ch.tag
                    if t != "extLst":
                        ti["type"] = t
                        break
            anims = []
            anis = set()
            tm = se.find(f"{{{PML}}}timing")
            if tm is not None:
                order = 0
                seen_targets: set = set()
                for par in tm.iter(f"{{{PML}}}par"):
                    # Skip non-leaf par nodes (containers that have child par
                    # elements) to avoid double-counting animation targets from
                    # ancestor containers.
                    if par.find(f"{{{PML}}}par") is not None:
                        continue
                    par_cTn = par.find(f"{{{PML}}}cTn")
                    trigger = "onClick"
                    delay_ms = 0
                    seq_dur = None
                    if par_cTn is not None:
                        nt = par_cTn.get("nodeType", "")
                        trigger = NODE_TYPE_TO_TRIGGER.get(nt, "unknown" if nt else "onClick")
                        d = par_cTn.get("delay", "0")
                        if d and d != "indefinite":
                            try:
                                delay_ms = int(d)
                            except Exception:
                                pass
                        dv = par_cTn.get("dur")
                        if dv and dv != "indefinite":
                            try:
                                seq_dur = int(dv)
                            except Exception:
                                pass
                    for tgt in par.findall(f".//{{{PML}}}spTgt"):
                        sp_id = tgt.get("spid")
                        # Deduplicate by (spid, par element id) to prevent
                        # the same target being counted more than once.
                        tgt_key = (sp_id, id(tgt))
                        if tgt_key in seen_targets:
                            continue
                        seen_targets.add(tgt_key)
                        sn = None
                        si = None
                        if sp_id:
                            for idx, sh in enumerate(slide.shapes):
                                try:
                                    if str(sh.shape_id) == str(sp_id):
                                        sn = sh.name
                                        si = idx
                                        anis.add(idx)
                                        break
                                except Exception:
                                    pass
                        ef_dur = seq_dur
                        for cb in par.findall(f".//{{{PML}}}cBhvr"):
                            ic = cb.find(f"{{{PML}}}cTn")
                            if ic is not None:
                                dv2 = ic.get("dur")
                                if dv2 and dv2 != "indefinite":
                                    try:
                                        ef_dur = int(dv2)
                                    except Exception:
                                        pass
                                break
                        et = "unknown"
                        for el in par.iter():
                            t = el.tag.split("}")[-1] if "}" in el.tag else el.tag
                            if t in ("animEffect", "anim", "animMotion", "animScale", "animRot", "set", "animClr"):
                                et = t
                                break
                        anims.append({
                            "order": order, "shape_name": sn, "shape_index": si,
                            "effect_type": et, "trigger": trigger,
                            "duration_ms": ef_dur, "delay_ms": delay_ms,
                        })
                        order += 1
            return {
                "session_id": session_id, "slide_index": slide_index,
                "has_animations": len(anims) > 0, "has_transition": has_t,
                "transition_info": ti, "animation_count": len(anims),
                "animations": anims, "animated_shape_indices": sorted(anis),
            }

    # ===== Batch1 新增工具 =====

    # shape_type 字符串到 MSO_SHAPE 枚举的映射
    _SHAPE_TYPE_MAP = {
        "RECTANGLE": MSO_SHAPE.RECTANGLE,
        "ROUNDED_RECTANGLE": MSO_SHAPE.ROUNDED_RECTANGLE,
        "OVAL": MSO_SHAPE.OVAL,
        "TRIANGLE": MSO_SHAPE.ISOSCELES_TRIANGLE,
        "ISOSCELES_TRIANGLE": MSO_SHAPE.ISOSCELES_TRIANGLE,
        "RIGHT_TRIANGLE": MSO_SHAPE.RIGHT_TRIANGLE,
        "DIAMOND": MSO_SHAPE.DIAMOND,
        "PENTAGON": MSO_SHAPE.PENTAGON,
        "HEXAGON": MSO_SHAPE.HEXAGON,
        "ARROW_RIGHT": MSO_SHAPE.RIGHT_ARROW,
        "RIGHT_ARROW": MSO_SHAPE.RIGHT_ARROW,
        "ARROW_LEFT": MSO_SHAPE.LEFT_ARROW,
        "LEFT_ARROW": MSO_SHAPE.LEFT_ARROW,
        "ARROW_UP": MSO_SHAPE.UP_ARROW,
        "UP_ARROW": MSO_SHAPE.UP_ARROW,
        "ARROW_DOWN": MSO_SHAPE.DOWN_ARROW,
        "DOWN_ARROW": MSO_SHAPE.DOWN_ARROW,
        "STAR_5_POINT": MSO_SHAPE.STAR_5_POINT,
        "HEART": MSO_SHAPE.HEART,
        "CLOUD": MSO_SHAPE.CLOUD,
    }

    # chart_type 字符串到 XL_CHART_TYPE 枚举的映射
    _CHART_TYPE_MAP = {
        "CLUSTERED_COLUMN": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "COLUMN_CLUSTERED": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "CLUSTERED_BAR": XL_CHART_TYPE.BAR_CLUSTERED,
        "BAR_CLUSTERED": XL_CHART_TYPE.BAR_CLUSTERED,
        "LINE": XL_CHART_TYPE.LINE,
        "PIE": XL_CHART_TYPE.PIE,
    }

    # alignment 字符串到 PP_ALIGN 枚举的映射
    _ALIGNMENT_MAP = {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
        "justify": PP_ALIGN.JUSTIFY,
    }

    @staticmethod
    def _parse_hex_color(hex_str: str) -> RGBColor:
        """将十六进制颜色字符串（如 'FF0000'）解析为 RGBColor。

        Args:
            hex_str: 6 位十六进制颜色字符串

        Returns:
            RGBColor 对象

        Raises:
            ValueError: 格式不合法
        """
        if not isinstance(hex_str, str):
            raise ValueError("颜色必须是字符串")
        cleaned = hex_str.strip().lstrip("#")
        if len(cleaned) != 6 or not all(c in "0123456789abcdefABCDEF" for c in cleaned):
            raise ValueError(f"无效的颜色格式: {hex_str!r}（需要 6 位十六进制如 'FF0000'）")
        r = int(cleaned[0:2], 16)
        g = int(cleaned[2:4], 16)
        b = int(cleaned[4:6], 16)
        return RGBColor(r, g, b)

    def add_shape(
        self,
        session_id: str,
        slide_index: int,
        shape_type: str,
        left: int,
        top: int,
        width: int,
        height: int,
        text: Optional[str] = None,
        fill_color: Optional[str] = None,
        line_color: Optional[str] = None,
    ) -> Dict[str, Any]:
        """添加 Auto Shape 到幻灯片。

        Args:
            session_id: 会话 ID
            slide_index: 幻灯片索引（0-based）
            shape_type: 形状类型字符串（如 RECTANGLE, OVAL, TRIANGLE 等）
            left: 左边距（EMU）
            top: 上边距（EMU）
            width: 宽度（EMU）
            height: 高度（EMU）
            text: 形状内文本（可选）
            fill_color: 填充颜色十六进制字符串如 'FF0000'（可选）
            line_color: 线条颜色十六进制字符串（可选）

        Returns:
            添加结果
        """
        session = self.sessions.get(session_id)

        # 参数校验
        shape_type_upper = shape_type.strip().upper() if isinstance(shape_type, str) else ""
        if shape_type_upper not in self._SHAPE_TYPE_MAP:
            supported = ", ".join(sorted(self._SHAPE_TYPE_MAP.keys()))
            raise ValueError(f"不支持的 shape_type: {shape_type!r}（支持: {supported}）")

        if not isinstance(left, (int, float)) or isinstance(left, bool):
            raise TypeError("left 必须是数字")
        if not isinstance(top, (int, float)) or isinstance(top, bool):
            raise TypeError("top 必须是数字")
        if not isinstance(width, (int, float)) or isinstance(width, bool):
            raise TypeError("width 必须是数字")
        if not isinstance(height, (int, float)) or isinstance(height, bool):
            raise TypeError("height 必须是数字")
        if int(width) <= 0:
            raise ValueError("width 必须大于 0")
        if int(height) <= 0:
            raise ValueError("height 必须大于 0")

        if text is not None:
            if not isinstance(text, str):
                raise TypeError("text 必须是字符串")
            if len(text) > limits.MAX_TEXT_LENGTH:
                raise ValueError(f"文本过长: {len(text)} > {limits.MAX_TEXT_LENGTH}")

        # 预解析颜色（锁外）
        parsed_fill = self._parse_hex_color(fill_color) if fill_color is not None else None
        parsed_line = self._parse_hex_color(line_color) if line_color is not None else None

        mso_shape = self._SHAPE_TYPE_MAP[shape_type_upper]

        with session.lock:
            prs = session.presentation
            slide_index = self._validate_slide_index(prs, slide_index)
            slide = prs.slides[slide_index]

            shape = slide.shapes.add_shape(
                mso_shape,
                Emu(int(left)),
                Emu(int(top)),
                Emu(int(width)),
                Emu(int(height)),
            )

            if text is not None:
                shape.text = text

            if parsed_fill is not None:
                shape.fill.solid()
                shape.fill.fore_color.rgb = parsed_fill

            if parsed_line is not None:
                shape.line.fill.solid()
                shape.line.color.rgb = parsed_line

            session.dirty = True

            return {
                "slide_index": slide_index,
                "shape_type": shape_type_upper,
                "shape_name": shape.name,
                "message": "形状已添加",
            }

    def add_chart(
        self,
        session_id: str,
        slide_index: int,
        chart_type: str,
        categories: List[str],
        series_data: Dict[str, List],
        left: int,
        top: int,
        width: int,
        height: int,
        title: Optional[str] = None,
    ) -> Dict[str, Any]:
        """添加图表到幻灯片。

        Args:
            session_id: 会话 ID
            slide_index: 幻灯片索引（0-based）
            chart_type: 图表类型字符串（CLUSTERED_COLUMN, CLUSTERED_BAR, LINE, PIE）
            categories: 分类标签列表
            series_data: 系列数据字典 {系列名称: [数据值]}
            left: 左边距（EMU）
            top: 上边距（EMU）
            width: 宽度（EMU）
            height: 高度（EMU）
            title: 图表标题（可选）

        Returns:
            添加结果
        """
        session = self.sessions.get(session_id)

        # 参数校验
        chart_type_upper = chart_type.strip().upper() if isinstance(chart_type, str) else ""
        if chart_type_upper not in self._CHART_TYPE_MAP:
            supported = ", ".join(sorted(self._CHART_TYPE_MAP.keys()))
            raise ValueError(f"不支持的 chart_type: {chart_type!r}（支持: {supported}）")

        if not isinstance(categories, list) or len(categories) == 0:
            raise ValueError("categories 必须是非空列表")
        for i, cat in enumerate(categories):
            if not isinstance(cat, str):
                raise TypeError(f"categories[{i}] 必须是字符串")

        if not isinstance(series_data, dict) or len(series_data) == 0:
            raise ValueError("series_data 必须是非空字典")

        cat_len = len(categories)
        for name, values in series_data.items():
            if not isinstance(name, str):
                raise TypeError("series_data 的键必须是字符串")
            if not isinstance(values, list):
                raise TypeError(f"series_data[{name!r}] 必须是列表")
            if len(values) != cat_len:
                raise ValueError(
                    f"series_data[{name!r}] 长度 {len(values)} 与 categories 长度 {cat_len} 不匹配"
                )
            for j, v in enumerate(values):
                if not isinstance(v, (int, float)) or isinstance(v, bool):
                    raise TypeError(f"series_data[{name!r}][{j}] 必须是数字")

        if not isinstance(left, (int, float)) or isinstance(left, bool):
            raise TypeError("left 必须是数字")
        if not isinstance(top, (int, float)) or isinstance(top, bool):
            raise TypeError("top 必须是数字")
        if not isinstance(width, (int, float)) or isinstance(width, bool):
            raise TypeError("width 必须是数字")
        if not isinstance(height, (int, float)) or isinstance(height, bool):
            raise TypeError("height 必须是数字")
        if int(width) <= 0:
            raise ValueError("width 必须大于 0")
        if int(height) <= 0:
            raise ValueError("height 必须大于 0")

        if title is not None and not isinstance(title, str):
            raise TypeError("title 必须是字符串")

        xl_chart_type = self._CHART_TYPE_MAP[chart_type_upper]

        chart_data = ChartData()
        chart_data.categories = categories
        for name, values in series_data.items():
            chart_data.add_series(name, values)

        with session.lock:
            prs = session.presentation
            slide_index = self._validate_slide_index(prs, slide_index)
            slide = prs.slides[slide_index]

            chart_frame = slide.shapes.add_chart(
                xl_chart_type,
                Emu(int(left)),
                Emu(int(top)),
                Emu(int(width)),
                Emu(int(height)),
                chart_data,
            )

            if title is not None:
                chart_frame.chart.has_title = True
                chart_frame.chart.chart_title.has_text_frame = True
                chart_frame.chart.chart_title.text_frame.text = title

            session.dirty = True

            return {
                "slide_index": slide_index,
                "chart_type": chart_type_upper,
                "categories_count": len(categories),
                "series_count": len(series_data),
                "has_title": title is not None,
                "message": "图表已添加",
            }

    def manage_text(
        self,
        session_id: str,
        operation: str,
        slide_index: Optional[int] = None,
        text: Optional[str] = None,
        left: Optional[int] = None,
        top: Optional[int] = None,
        width: Optional[int] = None,
        height: Optional[int] = None,
        font_size: Optional[int] = None,
        font_name: Optional[str] = None,
        bold: Optional[bool] = None,
        italic: Optional[bool] = None,
        color: Optional[str] = None,
        alignment: Optional[str] = None,
        shape_index: Optional[int] = None,
    ) -> Dict[str, Any]:
        """统一文本管理：添加文本框、格式化现有形状文本、提取文本。

        Args:
            session_id: 会话 ID
            operation: 操作类型 ('add', 'format', 'extract')
            slide_index: 幻灯片索引（extract 时可选，不传则提取所有）
            text: 文本内容（add 时必填）
            left, top, width, height: 位置和尺寸（EMU，add 时必填）
            font_size: 字号（add 默认 18）
            font_name: 字体名称（可选）
            bold: 是否加粗
            italic: 是否斜体
            color: 十六进制颜色字符串（可选）
            alignment: 对齐方式 ('left', 'center', 'right', 'justify')
            shape_index: 形状索引（format 时必填）

        Returns:
            操作结果
        """
        if operation not in ("add", "format", "extract"):
            raise ValueError(f"无效的 operation: {operation!r}（支持: add, format, extract）")

        session = self.sessions.get(session_id)

        if operation == "add":
            return self._manage_text_add(
                session, slide_index, text, left, top, width, height,
                font_size, font_name, bold, italic, color, alignment,
            )
        elif operation == "format":
            return self._manage_text_format(
                session, slide_index, shape_index,
                bold, italic, font_size, font_name, color, alignment,
            )
        else:  # extract
            return self._manage_text_extract(session, slide_index)

    def _manage_text_add(
        self, session, slide_index, text, left, top, width, height,
        font_size, font_name, bold, italic, color, alignment,
    ) -> Dict[str, Any]:
        """内部方法：添加文本框。"""
        if slide_index is None:
            raise ValueError("add 操作需要 slide_index")
        if text is None:
            raise ValueError("add 操作需要 text")
        if not isinstance(text, str):
            raise TypeError("text 必须是字符串")
        if len(text) > limits.MAX_TEXT_LENGTH:
            raise ValueError(f"文本过长: {len(text)} > {limits.MAX_TEXT_LENGTH}")
        if left is None or top is None or width is None or height is None:
            raise ValueError("add 操作需要 left, top, width, height")

        for name, val in [("left", left), ("top", top), ("width", width), ("height", height)]:
            if not isinstance(val, (int, float)) or isinstance(val, bool):
                raise TypeError(f"{name} 必须是数字")
        if int(width) <= 0:
            raise ValueError("width 必须大于 0")
        if int(height) <= 0:
            raise ValueError("height 必须大于 0")

        # add 操作使用默认值（用户未指定时）
        if font_size is None:
            font_size = 18
        if bold is None:
            bold = False
        if italic is None:
            italic = False
        if alignment is None:
            alignment = "left"

        if not isinstance(font_size, int) or isinstance(font_size, bool):
            raise TypeError("font_size 必须是整数")
        if font_size < 1 or font_size > self.MAX_FONT_SIZE:
            raise ValueError(f"font_size 必须在 1-{self.MAX_FONT_SIZE} 之间")

        parsed_color = self._parse_hex_color(color) if color is not None else None

        align_key = alignment.strip().lower() if isinstance(alignment, str) else ""
        if align_key not in self._ALIGNMENT_MAP:
            raise ValueError(f"无效的 alignment: {alignment!r}")
        pp_align = self._ALIGNMENT_MAP[align_key]

        with session.lock:
            prs = session.presentation
            slide_index = self._validate_slide_index(prs, slide_index)
            slide = prs.slides[slide_index]

            textbox = slide.shapes.add_textbox(
                Emu(int(left)),
                Emu(int(top)),
                Emu(int(width)),
                Emu(int(height)),
            )
            tf = textbox.text_frame
            tf.text = text

            for paragraph in tf.paragraphs:
                paragraph.alignment = pp_align
                for run in paragraph.runs:
                    run.font.size = Pt(font_size)
                    run.font.bold = bold
                    run.font.italic = italic
                    if font_name:
                        run.font.name = font_name
                    if parsed_color is not None:
                        run.font.color.rgb = parsed_color

            session.dirty = True

            return {
                "slide_index": slide_index,
                "operation": "add",
                "text_length": len(text),
                "message": "文本框已添加",
            }

    def _manage_text_format(
        self, session, slide_index, shape_index,
        bold, italic, font_size, font_name, color, alignment,
    ) -> Dict[str, Any]:
        """内部方法：格式化现有形状的文本。"""
        if slide_index is None:
            raise ValueError("format 操作需要 slide_index")
        if shape_index is None:
            raise ValueError("format 操作需要 shape_index")
        if not isinstance(shape_index, int) or isinstance(shape_index, bool):
            raise TypeError("shape_index 必须是整数")

        parsed_color = self._parse_hex_color(color) if color is not None else None
        pp_align = None
        if alignment is not None:
            align_key = alignment.strip().lower() if isinstance(alignment, str) else ""
            if align_key not in self._ALIGNMENT_MAP:
                raise ValueError(f"无效的 alignment: {alignment!r}")
            pp_align = self._ALIGNMENT_MAP[align_key]

        with session.lock:
            prs = session.presentation
            slide_index = self._validate_slide_index(prs, slide_index)
            slide = prs.slides[slide_index]

            if shape_index < 0 or shape_index >= len(slide.shapes):
                raise ValueError(
                    f"shape_index {shape_index} 超出范围 (有效范围: 0-{len(slide.shapes) - 1})"
                )

            shape = slide.shapes[shape_index]
            if not shape.has_text_frame:
                raise ValueError(f"shape[{shape_index}] 没有文本框")

            for paragraph in shape.text_frame.paragraphs:
                if pp_align is not None:
                    paragraph.alignment = pp_align
                for run in paragraph.runs:
                    if bold is not None:
                        run.font.bold = bold
                    if italic is not None:
                        run.font.italic = italic
                    if font_size is not None:
                        if isinstance(font_size, int) and not isinstance(font_size, bool):
                            run.font.size = Pt(font_size)
                    if font_name is not None:
                        run.font.name = font_name
                    if parsed_color is not None:
                        run.font.color.rgb = parsed_color

            session.dirty = True

            return {
                "slide_index": slide_index,
                "shape_index": shape_index,
                "operation": "format",
                "message": "文本格式已更新",
            }

    def _manage_text_extract(self, session, slide_index) -> Dict[str, Any]:
        """内部方法：提取文本。"""
        with session.lock:
            prs = session.presentation

            if slide_index is not None:
                slide_index = self._validate_slide_index(prs, slide_index)
                slides_to_check = [(slide_index, prs.slides[slide_index])]
            else:
                slides_to_check = list(enumerate(prs.slides))

            result_slides = []
            total_chars = 0
            for s_idx, slide in slides_to_check:
                texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        texts.append(shape.text)
                joined = "\n".join(texts)
                total_chars += len(joined)
                result_slides.append({
                    "slide_index": s_idx,
                    "text": joined,
                    "char_count": len(joined),
                })

            return {
                "operation": "extract",
                "total_slides": len(result_slides),
                "total_chars": total_chars,
                "slides": result_slides,
            }

    def format_table_cell(
        self,
        session_id: str,
        slide_index: int,
        shape_index: int,
        row: int,
        col: int,
        text: Optional[str] = None,
        font_size: Optional[int] = None,
        bold: Optional[bool] = None,
        fill_color: Optional[str] = None,
        alignment: Optional[str] = None,
    ) -> Dict[str, Any]:
        """格式化表格单元格。

        Args:
            session_id: 会话 ID
            slide_index: 幻灯片索引（0-based）
            shape_index: 形状索引（必须是表格类型）
            row: 行索引（0-based）
            col: 列索引（0-based）
            text: 新文本内容（可选）
            font_size: 字号（可选）
            bold: 是否加粗（可选）
            fill_color: 单元格填充颜色十六进制字符串（可选）
            alignment: 对齐方式（可选）

        Returns:
            格式化结果
        """
        session = self.sessions.get(session_id)

        # 参数校验
        if not isinstance(shape_index, int) or isinstance(shape_index, bool):
            raise TypeError("shape_index 必须是整数")
        if not isinstance(row, int) or isinstance(row, bool):
            raise TypeError("row 必须是整数")
        if not isinstance(col, int) or isinstance(col, bool):
            raise TypeError("col 必须是整数")
        if row < 0:
            raise ValueError("row 不能为负数")
        if col < 0:
            raise ValueError("col 不能为负数")

        if text is not None:
            if not isinstance(text, str):
                raise TypeError("text 必须是字符串")
            if len(text) > limits.MAX_TEXT_LENGTH:
                raise ValueError(f"文本过长: {len(text)} > {limits.MAX_TEXT_LENGTH}")

        if font_size is not None:
            if not isinstance(font_size, int) or isinstance(font_size, bool):
                raise TypeError("font_size 必须是整数")
            if font_size < 1 or font_size > self.MAX_FONT_SIZE:
                raise ValueError(f"font_size 必须在 1-{self.MAX_FONT_SIZE} 之间")

        parsed_fill = self._parse_hex_color(fill_color) if fill_color is not None else None

        pp_align = None
        if alignment is not None:
            align_key = alignment.strip().lower() if isinstance(alignment, str) else ""
            if align_key not in self._ALIGNMENT_MAP:
                raise ValueError(f"无效的 alignment: {alignment!r}")
            pp_align = self._ALIGNMENT_MAP[align_key]

        with session.lock:
            prs = session.presentation
            slide_index = self._validate_slide_index(prs, slide_index)
            slide = prs.slides[slide_index]

            if shape_index < 0 or shape_index >= len(slide.shapes):
                raise ValueError(
                    f"shape_index {shape_index} 超出范围 (有效范围: 0-{len(slide.shapes) - 1})"
                )

            shape = slide.shapes[shape_index]
            if not shape.has_table:
                raise ValueError(f"shape[{shape_index}] 不是表格")

            table = shape.table
            if row >= len(table.rows):
                raise ValueError(
                    f"row {row} 超出范围 (有效范围: 0-{len(table.rows) - 1})"
                )
            if col >= len(table.columns):
                raise ValueError(
                    f"col {col} 超出范围 (有效范围: 0-{len(table.columns) - 1})"
                )

            cell = table.cell(row, col)

            if text is not None:
                cell.text = text

            # 应用字体格式
            if font_size is not None or bold is not None or pp_align is not None:
                for paragraph in cell.text_frame.paragraphs:
                    if pp_align is not None:
                        paragraph.alignment = pp_align
                    for run in paragraph.runs:
                        if font_size is not None:
                            run.font.size = Pt(font_size)
                        if bold is not None:
                            run.font.bold = bold

            # 应用填充颜色
            if parsed_fill is not None:
                cell.fill.solid()
                cell.fill.fore_color.rgb = parsed_fill

            session.dirty = True

            return {
                "slide_index": slide_index,
                "shape_index": shape_index,
                "row": row,
                "col": col,
                "message": "表格单元格已格式化",
            }

    # ===== Batch2 新增工具 =====

    def manage_hyperlinks(
        self,
        session_id: str,
        slide_index: int,
        shape_index: int,
        operation: str,
        url: Optional[str] = None,
        text: Optional[str] = None,
    ) -> Dict[str, Any]:
        """管理形状中的超链接：添加、移除、列出、更新。

        Args:
            session_id: 会话 ID
            slide_index: 幻灯片索引（0-based）
            shape_index: 形状索引
            operation: 操作类型（add, remove, list, update）
            url: 超链接 URL（add/update 时必填）
            text: 超链接显示文本（add 时可选，指定时创建新 run）

        Returns:
            操作结果
        """
        valid_ops = ("add", "remove", "list", "update")
        if operation not in valid_ops:
            raise ValueError(
                f"无效的 operation: {operation!r}（支持: {', '.join(valid_ops)}）"
            )

        if operation in ("add", "update") and not url:
            raise ValueError(f"{operation} 操作需要提供 url 参数")

        if url is not None and not isinstance(url, str):
            raise TypeError("url 必须是字符串")
        if url is not None:
            _ALLOWED_SCHEMES = ("http://", "https://", "ftp://", "mailto:")
            if not any(url.lower().startswith(s) for s in _ALLOWED_SCHEMES):
                raise ValueError("URL 协议不受支持，请使用 http/https/ftp/mailto")
        if text is not None and not isinstance(text, str):
            raise TypeError("text 必须是字符串")

        if not isinstance(shape_index, int) or isinstance(shape_index, bool):
            raise TypeError("shape_index 必须是整数")

        session = self.sessions.get(session_id)

        with session.lock:
            prs = session.presentation
            slide_index = self._validate_slide_index(prs, slide_index)
            slide = prs.slides[slide_index]

            if shape_index < 0 or shape_index >= len(slide.shapes):
                raise ValueError(
                    f"shape_index {shape_index} 超出范围 (有效范围: 0-{len(slide.shapes) - 1})"
                )

            shape = slide.shapes[shape_index]
            if not shape.has_text_frame:
                raise ValueError(f"shape[{shape_index}] 没有文本框，无法操作超链接")

            text_frame = shape.text_frame

            if operation == "list":
                hyperlinks = []
                for para in text_frame.paragraphs:
                    for run in para.runs:
                        hl = run.hyperlink
                        if hl and hl.address:
                            hyperlinks.append({
                                "text": run.text,
                                "url": hl.address,
                            })
                return {
                    "slide_index": slide_index,
                    "shape_index": shape_index,
                    "operation": "list",
                    "hyperlinks": hyperlinks,
                    "count": len(hyperlinks),
                    "message": f"找到 {len(hyperlinks)} 个超链接",
                }

            elif operation == "add":
                if text:
                    # 创建新的 paragraph 和 run 携带超链接
                    p = text_frame.paragraphs[-1] if text_frame.paragraphs else text_frame.add_paragraph()
                    run = p.add_run()
                    run.text = text
                    run.hyperlink.address = url
                    added_count = 1
                else:
                    # 对所有现有 run 添加超链接
                    added_count = 0
                    for para in text_frame.paragraphs:
                        for run in para.runs:
                            run.hyperlink.address = url
                            added_count += 1

                session.dirty = True
                return {
                    "slide_index": slide_index,
                    "shape_index": shape_index,
                    "operation": "add",
                    "url": url,
                    "added_count": added_count,
                    "message": f"已为 {added_count} 个 run 添加超链接",
                }

            elif operation == "remove":
                removed_count = 0
                for para in text_frame.paragraphs:
                    for run in para.runs:
                        if run.hyperlink and run.hyperlink.address:
                            # 通过 lxml 删除超链接关系
                            r_elem = run._r
                            rPr = r_elem.find(qn("a:rPr"))
                            if rPr is not None:
                                hlinkClick = rPr.find(qn("a:hlinkClick"))
                                if hlinkClick is not None:
                                    rPr.remove(hlinkClick)
                                    removed_count += 1

                session.dirty = True
                return {
                    "slide_index": slide_index,
                    "shape_index": shape_index,
                    "operation": "remove",
                    "removed_count": removed_count,
                    "message": f"已移除 {removed_count} 个超链接",
                }

            else:  # update
                updated_count = 0
                for para in text_frame.paragraphs:
                    for run in para.runs:
                        if run.hyperlink and run.hyperlink.address:
                            run.hyperlink.address = url
                            updated_count += 1

                session.dirty = True
                return {
                    "slide_index": slide_index,
                    "shape_index": shape_index,
                    "operation": "update",
                    "url": url,
                    "updated_count": updated_count,
                    "message": f"已更新 {updated_count} 个超链接",
                }

    # 连接线类型映射（transition_type -> OOXML 元素名）
    _TRANSITION_TYPE_MAP = {
        "fade": "fade",
        "push": "push",
        "wipe": "wipe",
        "split": "split",
        "zoom": "zoom",
        "fly": "cover",
        "appear": "cut",      # appear 等效于 cut（无过渡效果时间=0）
        "dissolve": "dissolve",
        "cut": "cut",
        "wheel": "wheel",
        "strips": "strips",
        "checker": "checker",
        "blinds": "blinds",
        "box": "zoom",        # box 映射为 zoom
        "random": "random",
    }

    def add_connector(
        self,
        session_id: str,
        slide_index: int,
        start_x: float,
        start_y: float,
        end_x: float,
        end_y: float,
        line_color: Optional[str] = None,
        line_width: Optional[float] = None,
        arrow_start: bool = False,
        arrow_end: bool = True,
    ) -> Dict[str, Any]:
        """添加连接线/箭头到幻灯片。

        Args:
            session_id: 会话 ID
            slide_index: 幻灯片索引（0-based）
            start_x: 起点 X 坐标（厘米）
            start_y: 起点 Y 坐标（厘米）
            end_x: 终点 X 坐标（厘米）
            end_y: 终点 Y 坐标（厘米）
            line_color: 线条颜色十六进制字符串如 'FF0000'（可选）
            line_width: 线条宽度（磅，可选）
            arrow_start: 起点是否显示箭头（默认 False）
            arrow_end: 终点是否显示箭头（默认 True）

        Returns:
            添加结果
        """
        session = self.sessions.get(session_id)

        # 参数校验
        EMU_PER_CM = 360000

        for name, val in [("start_x", start_x), ("start_y", start_y),
                          ("end_x", end_x), ("end_y", end_y)]:
            if not isinstance(val, (int, float)) or isinstance(val, bool):
                raise TypeError(f"{name} 必须是数字")

        # 范围校验：坐标不能为负数
        start_x = self._validate_non_negative(start_x, "start_x")
        start_y = self._validate_non_negative(start_y, "start_y")
        end_x = self._validate_non_negative(end_x, "end_x")
        end_y = self._validate_non_negative(end_y, "end_y")

        if line_color is not None:
            parsed_color = self._parse_hex_color(line_color)
        else:
            parsed_color = None

        if line_width is not None:
            if not isinstance(line_width, (int, float)) or isinstance(line_width, bool):
                raise TypeError("line_width 必须是数字")
            if line_width <= 0:
                raise ValueError("line_width 必须大于 0")

        if not isinstance(arrow_start, bool):
            raise TypeError("arrow_start 必须是布尔值")
        if not isinstance(arrow_end, bool):
            raise TypeError("arrow_end 必须是布尔值")

        # 转换坐标为 EMU
        sx = int(start_x * EMU_PER_CM)
        sy = int(start_y * EMU_PER_CM)
        ex = int(end_x * EMU_PER_CM)
        ey = int(end_y * EMU_PER_CM)

        # 计算连接线的位置和尺寸
        left = min(sx, ex)
        top_pos = min(sy, ey)
        cx = abs(ex - sx)
        cy = abs(ey - sy)

        # 确保尺寸不为零（至少 1 EMU）
        if cx == 0:
            cx = 1
        if cy == 0:
            cy = 1

        # 确定是否需要翻转
        flip_h = "1" if ex < sx else "0"
        flip_v = "1" if ey < sy else "0"

        with session.lock:
            prs = session.presentation
            slide_index = self._validate_slide_index(prs, slide_index)
            slide = prs.slides[slide_index]

            # 通过 lxml 构建连接线 XML
            from lxml import etree

            # 获取唯一的 shape id
            existing_ids = set()
            for shape in slide.shapes:
                try:
                    existing_ids.add(shape.shape_id)
                except Exception:
                    pass
            new_id = max(existing_ids, default=0) + 1

            # 构建 cxnSp 元素
            cxnSp = etree.SubElement(
                slide._element.find(qn("p:cSld")).find(qn("p:spTree")),
                qn("p:cxnSp"),
            )

            # nvCxnSpPr
            nvCxnSpPr = etree.SubElement(cxnSp, qn("p:nvCxnSpPr"))
            cNvPr = etree.SubElement(nvCxnSpPr, qn("p:cNvPr"))
            cNvPr.set("id", str(new_id))
            cNvPr.set("name", f"Connector {new_id}")
            etree.SubElement(nvCxnSpPr, qn("p:cNvCxnSpPr"))
            etree.SubElement(nvCxnSpPr, qn("p:nvPr"))

            # spPr
            spPr = etree.SubElement(cxnSp, qn("p:spPr"))

            # xfrm
            xfrm = etree.SubElement(spPr, qn("a:xfrm"))
            if flip_h == "1":
                xfrm.set("flipH", "1")
            if flip_v == "1":
                xfrm.set("flipV", "1")

            off = etree.SubElement(xfrm, qn("a:off"))
            off.set("x", str(left))
            off.set("y", str(top_pos))
            ext = etree.SubElement(xfrm, qn("a:ext"))
            ext.set("cx", str(cx))
            ext.set("cy", str(cy))

            # prstGeom - 直线连接器
            prstGeom = etree.SubElement(spPr, qn("a:prstGeom"))
            prstGeom.set("prst", "straightConnector1")
            etree.SubElement(prstGeom, qn("a:avLst"))

            # 线条样式
            ln = etree.SubElement(spPr, qn("a:ln"))
            if line_width is not None:
                ln.set("w", str(int(line_width * 12700)))  # pt to EMU

            if parsed_color is not None:
                solidFill = etree.SubElement(ln, qn("a:solidFill"))
                srgbClr = etree.SubElement(solidFill, qn("a:srgbClr"))
                srgbClr.set("val", str(parsed_color))
            else:
                solidFill = etree.SubElement(ln, qn("a:solidFill"))
                srgbClr = etree.SubElement(solidFill, qn("a:srgbClr"))
                srgbClr.set("val", "000000")

            # 箭头
            if arrow_start:
                headEnd = etree.SubElement(ln, qn("a:headEnd"))
                headEnd.set("type", "triangle")
                headEnd.set("w", "med")
                headEnd.set("len", "med")

            if arrow_end:
                tailEnd = etree.SubElement(ln, qn("a:tailEnd"))
                tailEnd.set("type", "triangle")
                tailEnd.set("w", "med")
                tailEnd.set("len", "med")

            session.dirty = True

            return {
                "slide_index": slide_index,
                "start": {"x_cm": start_x, "y_cm": start_y},
                "end": {"x_cm": end_x, "y_cm": end_y},
                "arrow_start": arrow_start,
                "arrow_end": arrow_end,
                "message": "连接线已添加",
            }

    def manage_slide_transitions(
        self,
        session_id: str,
        slide_index: int,
        transition_type: str,
        duration: Optional[float] = None,
        advance_after: Optional[float] = None,
    ) -> Dict[str, Any]:
        """设置幻灯片过渡效果。

        Args:
            session_id: 会话 ID
            slide_index: 幻灯片索引（0-based）
            transition_type: 过渡类型（fade, push, wipe, split, zoom 等）
            duration: 过渡持续时间（秒，可选）
            advance_after: 自动切换时间（秒，可选）

        Returns:
            设置结果
        """
        session = self.sessions.get(session_id)

        # 参数校验
        tt_lower = transition_type.strip().lower() if isinstance(transition_type, str) else ""
        if tt_lower not in self._TRANSITION_TYPE_MAP:
            supported = ", ".join(sorted(self._TRANSITION_TYPE_MAP.keys()))
            raise ValueError(
                f"不支持的 transition_type: {transition_type!r}（支持: {supported}）"
            )

        if duration is not None:
            if not isinstance(duration, (int, float)) or isinstance(duration, bool):
                raise TypeError("duration 必须是数字")
            if duration < 0:
                raise ValueError("duration 不能为负数")

        if advance_after is not None:
            if not isinstance(advance_after, (int, float)) or isinstance(advance_after, bool):
                raise TypeError("advance_after 必须是数字")
            if advance_after < 0:
                raise ValueError("advance_after 不能为负数")

        ooxml_type = self._TRANSITION_TYPE_MAP[tt_lower]

        with session.lock:
            prs = session.presentation
            slide_index = self._validate_slide_index(prs, slide_index)
            slide = prs.slides[slide_index]

            se = slide._element
            PML = "http://schemas.openxmlformats.org/presentationml/2006/main"

            # 移除旧的 transition 节点
            old_trans = se.find(f"{{{PML}}}transition")
            if old_trans is not None:
                se.remove(old_trans)

            # 创建新的 transition 节点（必须在 timing/extLst 之前）
            from lxml import etree

            trans = etree.Element(qn("p:transition"))
            timing_el = se.find(qn('p:timing'))
            extlst_el = se.find(qn('p:extLst'))
            insert_before = timing_el if timing_el is not None else extlst_el
            if insert_before is not None:
                idx = list(se).index(insert_before)
                se.insert(idx, trans)
            else:
                se.append(trans)

            # 设置过渡持续时间（毫秒）
            if duration is not None:
                trans.set("spd", "med")  # 默认中速
                dur_ms = str(int(duration * 1000))
                trans.set("dur", dur_ms)

            # 设置自动切换时间
            if advance_after is not None:
                adv_ms = str(int(advance_after * 1000))
                trans.set("advTm", adv_ms)
                trans.set("advClick", "true")

            # 添加过渡类型子元素
            etree.SubElement(trans, qn(f"p:{ooxml_type}"))

            session.dirty = True

            return {
                "slide_index": slide_index,
                "transition_type": tt_lower,
                "ooxml_element": ooxml_type,
                "duration_seconds": duration,
                "advance_after_seconds": advance_after,
                "message": "幻灯片过渡效果已设置",
            }

    def set_core_properties(
        self,
        session_id: str,
        title: Optional[str] = None,
        subject: Optional[str] = None,
        author: Optional[str] = None,
        keywords: Optional[str] = None,
        comments: Optional[str] = None,
        category: Optional[str] = None,
    ) -> Dict[str, Any]:
        """设置文档核心属性。

        Args:
            session_id: 会话 ID
            title: 文档标题（可选）
            subject: 文档主题（可选）
            author: 作者（可选）
            keywords: 关键词（可选）
            comments: 备注（可选）
            category: 分类（可选）

        Returns:
            设置结果
        """
        session = self.sessions.get(session_id)

        # 参数校验
        fields = {
            "title": title,
            "subject": subject,
            "author": author,
            "keywords": keywords,
            "comments": comments,
            "category": category,
        }

        for name, val in fields.items():
            if val is not None and not isinstance(val, str):
                raise TypeError(f"{name} 必须是字符串")

        # 检查是否至少传入了一个参数
        if all(v is None for v in fields.values()):
            raise ValueError("至少需要设置一个属性（title, subject, author, keywords, comments, category）")

        with session.lock:
            prs = session.presentation
            props = prs.core_properties

            updated = {}
            if title is not None:
                props.title = title
                updated["title"] = title
            if subject is not None:
                props.subject = subject
                updated["subject"] = subject
            if author is not None:
                props.author = author
                updated["author"] = author
            if keywords is not None:
                props.keywords = keywords
                updated["keywords"] = keywords
            if comments is not None:
                props.comments = comments
                updated["comments"] = comments
            if category is not None:
                props.category = category
                updated["category"] = category

            session.dirty = True

            return {
                "updated_fields": updated,
                "field_count": len(updated),
                "message": f"已更新 {len(updated)} 个文档属性",
            }
